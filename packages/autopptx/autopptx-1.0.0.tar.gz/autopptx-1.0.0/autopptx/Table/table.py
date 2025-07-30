import logging
import textwrap
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.oxml import parse_xml
from pptx.util import Inches

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def validate_single_table(table):
    """
    Validate a single table as a non-empty 2D list with consistent row lengths.

    Parameters:
        table (list[list]): Table data to validate.

    Returns:
        bool: True if valid, False otherwise.
    """
    if not isinstance(table, list) or not table:
        return False
    if not all(isinstance(row, list) for row in table):
        return False
    row_len = len(table[0])
    if row_len == 0:
        return False
    for row in table:
        if len(row) != row_len:
            return False
    return True


def insert_table(shape, table_data):
    """
    Insert a new table into a blank table placeholder shape.

    Parameters:
        shape (pptx.shapes.placeholder.TablePlaceholder):
            A placeholder shape of type TABLE, currently empty.
        table_data (list[list[str]]):
            A 2D list representing the table content (rows × columns).

    Raises:
        ValueError: If the shape is not a table placeholder or data is invalid.
    """
    if not (
        shape.is_placeholder
        and shape.placeholder_format.type == PP_PLACEHOLDER.TABLE
    ):
        raise ValueError("The provided shape is not a table placeholder.")

    if not validate_single_table(table_data):
        raise ValueError("table_data must be a 2D list.")

    rows = len(table_data)
    cols = len(table_data[0])

    # Insert an empty table into the placeholder
    table_shape = shape.insert_table(rows, cols)
    table = table_shape.table

    # Fill in cell values
    for i in range(rows):
        for j in range(cols):
            table.cell(i, j).text = str(table_data[i][j])

    logger.info("✅ Inserted %d×%d table into placeholder.", rows, cols)


def replace_table(shape, table_data):
    """
    Replace an existing table (or insert one) in a table placeholder
    with new table content, while preserving the shape.

    If the placeholder does not contain a table yet (blank placeholder),
    this function will insert a new one. Otherwise, it rebuilds the table
    via XML.

    Parameters:
        shape (pptx.shapes.placeholder.TablePlaceholder):
            The table placeholder shape to modify.
        table_data (list[list[str]]):
            A 2D list representing the table content (rows × columns).

    Raises:
        ValueError: If the shape is not a table placeholder or data is invalid.
    """
    if not (
        shape.is_placeholder
        and shape.placeholder_format.type == PP_PLACEHOLDER.TABLE
    ):
        raise ValueError("The provided shape is not a table placeholder.")

    if not table_data or not all(isinstance(row, list) for row in table_data):
        raise ValueError("table_data must be a 2D list.")

    # Check if the shape already contains a table
    if not hasattr(shape, "table"):
        insert_table(shape, table_data)
        return

    rows = len(table_data)
    cols = max(len(row) for row in table_data)

    col_width = Inches(1.5)
    row_height = Inches(0.5)

    # Remove existing table structure from XML
    graphic_data = shape._element.graphic.graphicData
    for child in list(graphic_data):
        graphic_data.remove(child)

    # Build XML grid definition
    grid_cols_xml = "\n".join(f'<a:gridCol w="{int(col_width)}"/>' for _ in range(cols))

    # Build XML row and cell definitions
    rows_xml = ""
    for row in table_data:
        row_cells = ""
        for cell in row:
            cell_xml = textwrap.dedent(
                f"""
                <a:tc>
                  <a:txBody>
                    <a:bodyPr/>
                    <a:lstStyle/>
                    <a:p><a:r><a:t>{cell}</a:t></a:r></a:p>
                  </a:txBody>
                  <a:tcPr/>
                </a:tc>
                """
            )
            row_cells += cell_xml
        row_xml = f'<a:tr h="{int(row_height)}">{row_cells}</a:tr>'
        rows_xml += row_xml

    # Compose full table XML
    tbl_xml = textwrap.dedent(
        f"""
        <a:tbl xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
          <a:tblPr firstRow="1" bandRow="1"/>
          <a:tblGrid>
            {grid_cols_xml}
          </a:tblGrid>
          {rows_xml}
        </a:tbl>
        """
    )

    tbl_element = parse_xml(tbl_xml)
    graphic_data.append(tbl_element)

    logger.info("✅ Replaced table: %d rows × %d columns.", rows, cols)


def replace_table_cell(table_shape, row_idx, col_idx, new_text):
    """
    Replace the text content of a specific cell in a table placeholder.

    This function works on a table placeholder that already has a table inserted.
    It replaces the text in the specified cell (by row and column index).

    Parameters:
        table_shape (GraphicFrame): The table placeholder shape, must already
            contain a table.
        row_idx (int): Zero-based index of the target row.
        col_idx (int): Zero-based index of the target column.
        new_text (str): The new text to write into the cell.

    Raises:
        ValueError: If shape is not a table placeholder or does not contain a table.
        IndexError: If the specified cell index is out of bounds.
    """
    if not (
        table_shape.is_placeholder
        and table_shape.placeholder_format.type == PP_PLACEHOLDER.TABLE
    ):
        raise ValueError("The provided shape is not a table placeholder.")

    if not isinstance(new_text, str):
        raise ValueError("new_text must be a string.")

    if not hasattr(table_shape, "table"):
        raise ValueError("No table has been inserted into the placeholder yet.")

    table = table_shape.table

    if not (0 <= row_idx < len(table.rows)) or not (0 <= col_idx < len(table.columns)):
        raise IndexError(f"Cell index ({row_idx}, {col_idx}) is out of range.")

    cell = table.cell(row_idx, col_idx)
    cell.text = new_text

    logger.info('✅ Cell updated: row %d, column %d → "%s"', row_idx, col_idx, new_text)


if __name__ == "__main__":
    from pptx import Presentation

    prs = Presentation("./data/output_demo.pptx")
    slide = prs.slides[5]
    table_shape = None
    table = [
        ["Project", "Progress", "Owner"],
        ["A", "100%", "Alice"],
        ["B", "90%", "Bob"],
    ]

    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.TABLE:
            replace_table_cell(shape, 1, 2, "success ✅")
            table_shape = shape
            break

    replace_table(table_shape, table)
    prs.save("./data/output_table_cell.pptx")