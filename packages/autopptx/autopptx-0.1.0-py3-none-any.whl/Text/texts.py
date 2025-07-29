import logging

from autopptx.Type.find import find_placeholders

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def replace_title(slide, text):
    """
    Replace the content of the title placeholder.

    Parameters:
        slide (pptx.slide.Slide): The slide object to operate on.
        text (str): The title text to insert.
    """
    shapes = find_placeholders(slide, "title")
    if not shapes:
        logger.warning("⚠️ Title placeholder not found.")
        return

    shape = shapes[0]
    tf = shape.text_frame
    tf.clear()
    tf.paragraphs[0].text = str(text)


def replace_subtitle(slide, text):
    """
    Replace the content of the subtitle placeholder.

    Parameters:
        slide (pptx.slide.Slide): The slide object to operate on.
        text (str): The subtitle text to insert.
    """
    shapes = find_placeholders(slide, "subtitle")
    if not shapes:
        logger.warning("⚠️ Subtitle placeholder not found.")
        return

    shape = shapes[0]
    tf = shape.text_frame
    tf.clear()
    tf.paragraphs[0].text = str(text)


def replace_bodytexts(slide, text, distribute_to_multiple_boxes=True):
    """
    Replace body text content in the slide. Optionally distribute content
    across multiple placeholders.

    Parameters:
        slide (pptx.slide.Slide): The slide object to operate on.
        text (list[str] or str): Paragraph content to insert.
        distribute_to_multiple_boxes (bool):
            - If True: Assign each paragraph to a separate text box.
            - If False: Insert all paragraphs into a single text box.
    """
    if not isinstance(text, list):
        text = [str(text)]

    shapes = find_placeholders(slide, "bodytext")
    if not shapes:
        logger.warning("⚠️ Body text placeholder not found.")
        return

    if distribute_to_multiple_boxes:
        num_shape = len(shapes)
        num_text = len(text)
        if num_text < num_shape:
            logger.warning(
                "⚠️ Provided %d paragraphs for %d placeholders, only the first few will be used.",
                num_text,
                num_shape,
            )
        elif num_text > num_shape:
            logger.warning(
                "⚠️ Provided %d paragraphs for %d placeholders, excess text will be ignored.",
                num_text,
                num_shape,
            )

        for i, shape in enumerate(shapes):
            if i >= len(text):
                break
            tf = shape.text_frame
            tf.clear()
            tf.paragraphs[0].text = text[i]
    else:
        shape = shapes[0]
        if not shape.has_text_frame:
            return
        tf = shape.text_frame
        tf.clear()
        for i, para in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = para

    logger.info("✅ Replaced texts successfully")


if __name__ == "__main__":
    from pptx import Presentation

    prs = Presentation("./data/template.pptx")
    slide = prs.slides[1]

    data = {
        "title": "2025 Annual Summary",
        "subtitle": "Marketing Department",
        "bodytext": ["Part 1: Work Review", "Part 2: Planning Ahead"],
        "image": "images/chart1.png",
        "table": [
            [
                ["Project", "Progress", "Owner"],
                ["A", "100%", "Alice"],
                ["B", "90%", "Bob"]
            ]
        ],
    }

    # Single text box with multiple paragraphs
    replace_title(slide, data["title"])
    replace_bodytexts(slide, data["bodytext"], distribute_to_multiple_boxes=False)
    prs.save("./data/output1.pptx")

    # Each paragraph in a separate text box
    replace_bodytexts(slide, data["bodytext"], distribute_to_multiple_boxes=True)
    prs.save("./data/output2.pptx")