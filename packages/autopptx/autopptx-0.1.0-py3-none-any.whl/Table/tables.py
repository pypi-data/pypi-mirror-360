import logging

from autopptx.Type.find import find_placeholders
from autopptx.Table.table import replace_table,validate_single_table

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def validate_table_data(table_data):
    """
    Validate that table_data is a list of valid tables.

    Parameters:
        table_data (list): List of tables to validate.

    Returns:
        bool: True if all tables are valid, False otherwise.
    """
    if not isinstance(table_data, list):
        return False
    for table in table_data:
        if not validate_single_table(table):
            return False
    return True


def replace_tables(slide, table_list):
    """
    Replace all table placeholders in a slide using the provided list of table data.

    Each placeholder is replaced with a corresponding 2D table from `table_list`.
    The number of replacements is determined by the shorter of the two lists.

    Parameters:
        slide (pptx.slide.Slide): The slide containing table placeholders.
        table_list (list[list[list[str]]]): A list of 2D tables (list of rows),
            where each row is a list of strings.

    Notes:
        - If the number of provided tables is fewer than the number of placeholders,
          only the first few placeholders will be filled.
        - If the number of tables exceeds the number of placeholders, the extra
          tables will be ignored.
        - Each table must be a 2D list, and each row must be a list of strings.
    """
    shapes = find_placeholders(slide, "table")
    if not shapes:
        logger.warning("⚠️ No table placeholder found.")
        return

    if not validate_table_data(table_list):
        raise ValueError("Invalid table_list: must be a list of valid 2D tables")

    num_shapes = len(shapes)
    num_tables = len(table_list)

    if num_tables < num_shapes:
        logger.warning(
            "⚠️ Provided tables (%d) < placeholders (%d); only part of placeholders will be filled.",
            num_tables,
            num_shapes,
        )
    elif num_tables > num_shapes:
        logger.warning(
            "⚠️ Provided tables (%d) > placeholders (%d); extra tables will be ignored.",
            num_tables,
            num_shapes,
        )

    replaced = 0
    for i, shape in enumerate(shapes):
        if i >= num_tables:
            break

        table_data = table_list[i]
        if not table_data or not all(isinstance(row, list) for row in table_data):
            logger.error("❌ Invalid table format at index %d, skipped.", i)
            continue

        try:
            replace_table(shape, table_data)
            replaced += 1
        except Exception as e:
            logger.error("❌ Failed to replace table at index %d: %s", i, e)
            raise

    logger.info(
        "✅ Table placeholder replacement complete: %d / %d tables replaced.",
        replaced,
        min(num_shapes, num_tables),
    )


if __name__ == "__main__":
    from pptx import Presentation

    prs = Presentation("./data/template.pptx")
    slide = prs.slides[5]
    table_list = [
        [
            ["Project", "Progress", "Owner"],
            ["A", "100%", "Alice"],
            ["B", "90%", "Bob"],
        ],
        [
            ["Task", "Status", "Responsible"],
            ["Design", "Completed", "Charlie"],
            ["Implementation", "In Progress", "Dana"],
            ["Testing", "Pending", "Eve"],
        ],
    ]

    replace_tables(slide, table_list)
    prs.save("./data/output_tables.pptx")