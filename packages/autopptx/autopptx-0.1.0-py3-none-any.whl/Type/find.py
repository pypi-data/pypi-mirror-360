import logging
from pptx.enum.shapes import PP_PLACEHOLDER

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def find_placeholders(slide, placeholder_type):
    """
    Return a list of shapes matching the specified placeholder type
    in the given slide.

    Supported placeholder_type values include:
        - "title": Title placeholder, identified by
          PP_PLACEHOLDER.TITLE or slide.shapes.title
        - "subtitle": Subtitle placeholder, usually the text placeholder
          with idx = 1
        - "bodytext": Body text placeholders, text placeholders with idx >= 2
        - "image": Picture placeholders, type PP_PLACEHOLDER.PICTURE
        - "table": Table placeholders, type PP_PLACEHOLDER.TABLE

    Parameters:
        slide (pptx.slide.Slide): The slide object to search.
        placeholder_type (str): The type of placeholder to find.

    Returns:
        List[Shape]: A list of placeholder shapes matching the specified
            type. Returns an empty list if none found.

    Raises:
        ValueError: If an unsupported placeholder_type is provided.
    """
    shapes = []

    if placeholder_type == "title":
        if slide.shapes.title:
            shapes.append(slide.shapes.title)

    elif placeholder_type == "subtitle":
        for shape in slide.placeholders:
            if (
                shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE
                and shape.has_text_frame
            ):
                shapes.append(shape)

    elif placeholder_type == "bodytext":
        for shape in slide.placeholders:
            if (
                shape.placeholder_format.type == PP_PLACEHOLDER.BODY
                and shape.has_text_frame
            ):
                shapes.append(shape)

    elif placeholder_type == "image":
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                shapes.append(shape)

    elif placeholder_type == "table":
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TABLE:
                shapes.append(shape)

    else:
        raise ValueError(f"Unsupported placeholder_type: {placeholder_type}")

    logger.debug(
        "Found %d placeholder(s) for type '%s' in slide %d",
        len(shapes),
        placeholder_type,
        slide.slide_id,
    )
    return shapes