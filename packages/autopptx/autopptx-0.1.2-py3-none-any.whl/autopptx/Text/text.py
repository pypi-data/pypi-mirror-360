import logging
from pptx.enum.shapes import PP_PLACEHOLDER

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def replace_bodytext(shape, text):
    """
    Replace the content of a body text placeholder with the specified text.

    This function clears the text frame of the placeholder and inserts new content.
    It raises an error if the given shape is not a valid body text placeholder.

    Parameters:
        shape (pptx.shapes.placeholder.Placeholder):
            The placeholder shape to update.
        text (str):
            The new body text content to insert.

    Raises:
        ValueError: If the provided shape is not a body text placeholder.
    """
    if not (
        shape.is_placeholder
        and shape.placeholder_format.type == PP_PLACEHOLDER.BODY
    ):
        raise ValueError("The provided shape is not a body text placeholder.")

    try:
        tf = shape.text_frame
        tf.clear()
        tf.paragraphs[0].text = text
        logger.info(
            "✅ Replaced text: %s -> placeholder %s",
            text,
            shape.placeholder_format.idx,
        )
    except Exception as e:
        logger.error(
            "❌ Failed to replace text in placeholder %s: %s",
            shape.placeholder_format.idx,
            e,
        )
        raise


if __name__ == "__main__":
    from pptx import Presentation
    import sys

    prs = Presentation("./data/template.pptx")
    slide = prs.slides[0]
    bodytext_shape = None
    text = "Part 1: Work Review"

    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.BODY:
            bodytext_shape = shape
            break

    if bodytext_shape is None:
        logger.error("No body text placeholder found on the first slide.")
        sys.exit(1)

    replace_bodytext(bodytext_shape, text)
    prs.save("./data/output_text.pptx")