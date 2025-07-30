import logging
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)

# Text alignment map
ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
    "distribute": PP_ALIGN.DISTRIBUTE,
}


def set_table_cell(
    cell,
    font_name="Arial",
    font_size=18,
    bold=False,
    italic=False,
    font_color=(0, 0, 0),
    align="center",
    bg_color=None
):
    """
    Apply text and background style to a table cell.

    Parameters:
        cell (pptx.table._Cell): Table cell
        font_name (str): Font family
        font_size (int): Font size in pt
        bold (bool): Bold text
        italic (bool): Italic text
        font_color (tuple): RGB font color
        align (str): Text alignment
        bg_color (tuple): RGB background color
    """
    if not hasattr(cell, "text_frame") or not cell.text_frame:
        raise ValueError("Provided cell has no text_frame, cannot style.")

    if align.lower() not in ALIGN_MAP:
        raise ValueError(
            f"Invalid alignment: '{align}'. Must be one of {list(ALIGN_MAP.keys())}"
        )

    for paragraph in cell.text_frame.paragraphs:
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.bold = bold
        font.italic = italic
        font.color.rgb = RGBColor(*font_color)

        if align.lower() in ALIGN_MAP:
            paragraph.alignment = ALIGN_MAP[align.lower()]

    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*bg_color)


def set_table_style(
    shape,
    font_name="Arial",
    font_size=18,
    bold=False,
    italic=False,
    font_color=(0, 0, 0),
    align="center",
    bg_color=None
):
    """
    Apply style to all cells in a shape that contains a table.
    """
    if not hasattr(shape, "table") or shape.table is None:
        logger.warning(f"⛔️ Shape '{shape.name}' does not contain a table. Skipping.")
        return
    
    table = shape.table
    logger.info("🎯 Applying table style...")

    for row in table.rows:
        for cell in row.cells:
            set_table_cell(
                cell,
                font_name=font_name,
                font_size=font_size,
                bold=bold,
                italic=italic,
                font_color=font_color,
                align=align,
                bg_color=bg_color
            )
    logger.info("✅ Table styling complete!")


def _get_rgb_safe(color_obj):
    try:
        if color_obj and color_obj.rgb:
            return (color_obj.rgb[0], color_obj.rgb[1], color_obj.rgb[2])
    except AttributeError:
        return None
    return None


def extract_cell_style(cell):
    """
    Extract style information from a table cell.

    Returns:
        dict: style properties like font, color, align, etc.
    """
    style = {}

    if not cell.text_frame or not cell.text_frame.paragraphs:
        logger.warning("⚠️ Empty cell, returning empty style.")
        return style

    para = cell.text_frame.paragraphs[0]
    run = para.runs[0] if para.runs else para.add_run()
    font = run.font

    if font:
        if font.name:
            style["font_name"] = font.name
        if font.size:
            style["font_size"] = font.size.pt
        if font.bold is not None:
            style["bold"] = font.bold
        if font.italic is not None:
            style["italic"] = font.italic
        rgb = _get_rgb_safe(font.color)
        if rgb:
            style["font_color"] = rgb

    if para.alignment:
        style["align"] = para.alignment.name

    try:
        cell.fill.solid()  # Make sure the fore_color is accessible
        bg_rgb = _get_rgb_safe(cell.fill.fore_color)
        if bg_rgb:
            style["bg_color"] = bg_rgb
    except Exception:
        pass

    # logger.info("✅ Cell style extraction complete!")
    return style


def extract_table_style(shape):
    """
    Extract styles from an entire table shape.

    Parameters:
        shape (pptx.shapes.Shape): A shape that contains a table.

    Returns:
        list: 2D list of style dicts
    """
    if not hasattr(shape, "table") or shape.table is None:
        logger.warning(f"⛔️ Shape '{shape.name}' does not contain a table. Skipping.")
        return []
    
    table = shape.table
    rows, cols = table.rows.__len__(), table.columns.__len__()
    logger.info(f"🎯 Extracting styles from table: {rows} rows x {cols} cols")

    style_table = []
    for row in table.rows:
        row_styles = []
        for cell in row.cells:
            row_styles.append(extract_cell_style(cell))
        style_table.append(row_styles)

    logger.info("✅ Table style extraction complete!")
    return style_table


def transfer_table_style(src_shape, dst_shape, mode="full"):
    """
    Transfer table styles from source shape to destination shape.

    Parameters:
        src_shape: Source shape with table
        dst_shape: Destination shape with table
        mode (str): 'full' for per-cell transfer, 'single' for uniform style
    """
    src_table = src_shape.table
    dst_table = dst_shape.table

    if mode == "single":
        base_style = extract_cell_style(src_table.cell(0, 0))
        if not base_style:
            logger.warning("⚠️ Source (0,0) cell style empty, skipping transfer.")
            return

        dst_rows = len(dst_table.rows)
        dst_cols = len(dst_table.columns)

        for i in range(dst_rows):
            for j in range(dst_cols):
                set_table_cell(dst_table.cell(i, j), **base_style)

    elif mode == "full":
        rows = min(len(src_table.rows), len(dst_table.rows))
        cols = min(len(src_table.columns), len(dst_table.columns))

        logger.info("📐 Applying styles cell-by-cell.")
        for i in range(rows):
            for j in range(cols):
                src_cell = src_table.cell(i, j)
                dst_cell = dst_table.cell(i, j)

                style = extract_cell_style(src_cell)
                if style:
                    set_table_cell(dst_cell, **style)
                else:
                    logger.warning(f"⚠️ No style found for row {i+1}, col {j+1}, skipping")

    else:
        raise ValueError("mode must be either 'full' or 'single'")

    logger.info("✅ Table style transfer complete!")


if __name__ == '__main__':
    from pptx import Presentation

    prs = Presentation("data/output_demo.pptx")
    table1 = prs.slides[5].shapes[2]
    table2 = prs.slides[5].shapes[3]

    set_table_style(
        table1,
        font_name="Microsoft YaHei",
        font_size=18,
        bold=True,
        italic=False,
        font_color=(255, 255, 255),
        align="center",
        bg_color=(30, 144, 255)
    )

    logger.info(extract_table_style(table1))
    prs.save("data/table_styled_full.pptx")

    transfer_table_style(table1, table2)
    prs.save("data/table_styled_transfer_full.pptx")

    transfer_table_style(table1, table2, mode="single")
    prs.save("data/table_styled_transfer_single.pptx")