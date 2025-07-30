import os
import logging
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.parts.image import Image

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def replace_image(shape, image_path):
    """
    Replace or insert the image content of a picture placeholder without removing
    the shape itself.

    This function handles two scenarios:
    1. If the placeholder is empty (no image yet), it inserts the new image.
    2. If the placeholder already contains an image, it replaces the underlying
       image data (blob override) to update the picture efficiently.

    Parameters:
        shape (pptx.shapes.placeholder.PicturePlaceholder or None):
            The picture placeholder shape to update. If None, the function returns early.
        image_path (str): Absolute or relative path to the new image file.

    Raises:
        ValueError: If the provided shape is not a picture placeholder.

    Logs:
        Info and error messages on success or failure of insertion or replacement.
    """
    if shape is None:
        logger.error("No shape provided to replace_image() — shape is None.")
        return

    if not (
        shape.is_placeholder
        and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE
    ):
        raise ValueError("The provided shape is not a picture placeholder.")

    img_path = os.path.abspath(image_path)

    # Check if image file exists
    if not os.path.exists(img_path):
        raise FileNotFoundError(f"Image file not found: {img_path}")

    # Case 1: Placeholder without an image (no blip_rId yet)
    if not hasattr(shape._element, "blip_rId"):
        try:
            shape.insert_picture(img_path)
            logger.info("✅ Inserted image into empty placeholder: %s", img_path)
        except Exception as e:
            logger.error(
                "❌ Failed to insert image into placeholder %d: %s",
                shape.placeholder_format.idx,
                e,
            )
            raise
    # Case 2: Placeholder with existing image (replace blob content)
    else:
        try:
            slide_part = shape.part
            blip_rId = shape._element.blip_rId
            img_part = slide_part.related_part(blip_rId)

            new_img = Image.from_file(img_path)
            img_part.blob = new_img._blob

            logger.info(
                "✅ Replaced image: %s -> placeholder %d",
                img_path,
                shape.placeholder_format.idx,
            )
        except Exception as e:
            logger.error(
                "❌ Failed to replace image in placeholder %d: %s",
                shape.placeholder_format.idx,
                e,
            )
            raise


if __name__ == "__main__":
    from pptx import Presentation

    prs = Presentation("./data/template.pptx")
    slide = prs.slides[2]
    picture_shape = None
    image_path = "./data/cat1.png"

    for shape in slide.placeholders:
        if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
            picture_shape = shape
            break

    replace_image(picture_shape, image_path)
    prs.save("./data/output_image.pptx")