import os
import logging

from autopptx.Type.find import find_placeholders
from autopptx.Image.image import replace_image

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def replace_images(slide, image_path):
    """
    Replace the image content of picture placeholders on the slide,
    including placeholders that don't yet contain images.

    Parameters:
        slide (pptx.slide.Slide): The slide object containing picture placeholders.
        image_path (str or list[str]): Path(s) to new image(s). If a list,
            each image will replace one placeholder in order.
    """
    if not isinstance(image_path, list):
        image_path = [str(image_path)]

    shapes = find_placeholders(slide, "image")
    if not shapes:
        logger.warning("⚠️ Picture placeholder not found.")
        return

    num_shapes = len(shapes)
    num_images = len(image_path)

    if num_images < num_shapes:
        logger.warning(
            "⚠️ The number of provided images (%d) is less than "
            "the number of placeholders (%d); only the first ones will be replaced.",
            num_images,
            num_shapes,
        )
    elif num_images > num_shapes:
        logger.warning(
            "⚠️ The number of provided images (%d) exceeds "
            "the number of placeholders (%d); extra images will be ignored.",
            num_images,
            num_shapes,
        )

    for i, shape in enumerate(shapes):
        if i >= num_images:
            break  # Prevent index out of range

        img_path = os.path.abspath(image_path[i])
        try:
            replace_image(shape, img_path)
            logger.info("Replaced image in placeholder %d with %s", i, img_path)
        except Exception as e:
            logger.error(
                "Failed to replace image in placeholder %d with %s: %s",
                i,
                img_path,
                e,
            )
            raise


if __name__ == "__main__":
    from pptx import Presentation

    prs = Presentation("./data/template.pptx")
    slide = prs.slides[3]
    image_path = ["./data/bunny1.png", "./data/bunny2.png"]

    replace_images(slide, image_path)
    prs.save("./data/output_images.pptx")