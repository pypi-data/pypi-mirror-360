import logging
from pptx.enum.shapes import PP_PLACEHOLDER

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def is_text(shape):
    """
    Check whether the given shape has a text frame.

    Parameters:
        shape (Shape): The shape object to check.

    Returns:
        bool: True if the shape contains a text frame, False otherwise.
    """
    return shape.has_text_frame


def is_title(shape):
    """
    Check whether the given shape is a title placeholder.

    A title placeholder typically has type TITLE or CENTER_TITLE.

    Parameters:
        shape (Shape): The shape object to check.

    Returns:
        bool: True if the shape is a title placeholder, False otherwise.
    """
    return (
        shape.is_placeholder
        and shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    )


def is_subtitle(shape):
    """
    Check whether the given shape is a subtitle placeholder.

    Parameters:
        shape (Shape): The shape object to check.

    Returns:
        bool: True if the shape is a subtitle placeholder, False otherwise.
    """
    return shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE


def is_bodytext(shape):
    """
    Check whether the given shape is a body text placeholder.

    Parameters:
        shape (Shape): The shape object to check.

    Returns:
        bool: True if the shape is a body text placeholder, False otherwise.
    """
    return shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.BODY


def is_image(shape):
    """
    Check whether the given shape is a picture placeholder.

    Parameters:
        shape (Shape): The shape object to check.

    Returns:
        bool: True if the shape is a picture placeholder, False otherwise.
    """
    return shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE


def is_table(shape):
    """
    Check whether the given shape is a table placeholder.

    Parameters:
        shape (Shape): The shape object to check.

    Returns:
        bool: True if the shape is a table placeholder, False otherwise.
    """
    return shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TABLE


if __name__ == "__main__":
    from pptx import Presentation

    prs = Presentation("./data/output_demo.pptx")
    slide = prs.slides[4]

    for idx, shape in enumerate(slide.shapes):
        logger.info(f"\nShape {idx}:")
        logger.info(f"- Is text box: {is_text(shape)}")
        logger.info(f"- Is title: {is_title(shape)}")
        logger.info(f"- Is subtitle: {is_subtitle(shape)}")
        logger.info(f"- Is body text: {is_bodytext(shape)}")
        logger.info(f"- Is picture placeholder: {is_image(shape)}")
        logger.info(f"- Is table placeholder: {is_table(shape)}")

        if shape.is_placeholder:
            phf = shape.placeholder_format
            logger.info(f"- Placeholder index: {phf.idx}")
            logger.info(f"- Placeholder type: {phf.type}")
        else:
            logger.info(f"- Non-placeholder shape")