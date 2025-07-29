import logging
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)

# Text alignment mapping
ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
    "distribute": PP_ALIGN.DISTRIBUTE,
}


def set_paragraph_style(
    paragraph,
    font_name="Arial",
    font_size=18,
    bold=False,
    italic=False,
    font_color=(0, 0, 0),
    align="left",
):
    """
    Set font style and alignment for a single paragraph.

    Parameters:
        paragraph (pptx.text.Paragraph): The paragraph to style.
        font_name (str): Font family name.
        font_size (int): Font size in points.
        bold (bool): Whether to apply bold.
        italic (bool): Whether to apply italic.
        font_color (tuple): RGB tuple for font color.
        align (str): Text alignment ("left", "center", etc.).
    """
    if not hasattr(paragraph, "runs") or not hasattr(paragraph, "alignment"):
        raise ValueError("Input does not appear to be a paragraph object.")
    
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    font = run.font

    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*font_color)

    align_value = ALIGN_MAP.get(align.lower())
    if align_value is not None:
        paragraph.alignment = align_value
        logger.debug(
            f"Applied style: font={font_name}, size={font_size}, "
            f"bold={bold}, italic={italic}, color={font_color}, align={align}"
        )
    else:
        raise ValueError(
            f"Unsupported align value: '{align}'. Must be one of {list(ALIGN_MAP.keys())}"
        )


def set_textbox_style(
    shape,
    font_name="Arial",
    font_size=18,
    bold=False,
    italic=False,
    font_color=(0, 0, 0),
    align="left",
):
    """
    Apply text style to all paragraphs in a text box.

    Parameters:
        shape (pptx.shapes.base.Shape): The shape with a text frame.
        font_name (str): Font family name.
        font_size (int): Font size in points.
        bold (bool): Whether to apply bold.
        italic (bool): Whether to apply italic.
        font_color (tuple): RGB tuple for font color.
        align (str): Text alignment ("left", "center", "right", etc.).
    """
    if not shape.has_text_frame:
        logger.warning("This shape has no text frame.")
        return

    for paragraph in shape.text_frame.paragraphs:
        set_paragraph_style(
            paragraph,
            font_name=font_name,
            font_size=font_size,
            bold=bold,
            italic=italic,
            font_color=font_color,
            align=align,
        )
    
    logger.info("✅ Text styling complete!")


def extract_paragraph_style(paragraph):
    """
    Extract text style attributes from the first run of a paragraph.

    Returns:
        dict: A dictionary with font name, size (pt), color (RGB tuple),
              bold/italic flags, and paragraph alignment.
    """
    if not hasattr(paragraph, "runs") or not hasattr(paragraph, "alignment"):
        raise ValueError("Input does not appear to be a paragraph object.")
    
    run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
    font = run.font

    style = {
        "font_name": font.name if font.name else None,
        "font_size": font.size.pt if font.size else None,
        "bold": font.bold,
        "italic": font.italic,
        "font_color": None,
        "align": paragraph.alignment.name if paragraph.alignment else None,
    }

    if font.color and font.color.rgb:
        rgb = font.color.rgb
        style["font_color"] = (rgb[0], rgb[1], rgb[2])

    return style


def extract_textbox_style(shape):
    """
    Extract style info from all paragraphs in a text box shape.

    Parameters:
        shape (pptx.shapes.base.Shape): The shape with a text frame.

    Returns:
        list[dict]: A list of paragraph style dictionaries.
    """
    if not shape.has_text_frame or not shape.text_frame.paragraphs:
        logger.warning("⚠️ The shape has no text content. Cannot extract styles.")
        return []

    style_list = []
    for para in shape.text_frame.paragraphs:
        style = extract_paragraph_style(para)
        style_list.append(style)

    logger.info("✅ Text style extraction complete!")
    return style_list


def transfer_text_style(src_shape, dst_shape, mode="full"):
    """
    Transfer text styles from source shape to destination shape.

    Parameters:
        src_shape: Source shape with text_frame
        dst_shape: Destination shape with text_frame
        mode (str): 'full' for per-paragraph transfer, 'single' for uniform style
    """
    if not (src_shape.has_text_frame and dst_shape.has_text_frame):
        logger.warning("⚠️ Either source or destination has no text frame. Skipping.")
        return

    src_paragraphs = src_shape.text_frame.paragraphs
    dst_paragraphs = dst_shape.text_frame.paragraphs

    if mode == "single":
        base_style = extract_paragraph_style(src_paragraphs[0])
        if not base_style:
            logger.warning("⚠️ Source first paragraph has no style, skipping transfer.")
            return

        logger.info("🎯 Applying style of first paragraph to all destination paragraphs.")
        for dst_para in dst_paragraphs:
            set_paragraph_style(dst_para, **base_style)

    elif mode == "full":
        count = min(len(src_paragraphs), len(dst_paragraphs))
        logger.info("📐 Applying paragraph styles one by one.")
        for i in range(count):
            src_para = src_paragraphs[i]
            dst_para = dst_paragraphs[i]

            style = extract_paragraph_style(src_para)
            if style:
                set_paragraph_style(dst_para, **style)
            else:
                logger.warning(f"⚠️ No style found for paragraph {i + 1}, skipping")

    else:
        raise ValueError("mode must be either 'full' or 'single'")

    logger.info("✅ Text style transfer complete!")


if __name__ == "__main__":
    from pptx import Presentation
    from autopptx.Type.find import find_placeholders

    pptx_file = "data/output_demo.pptx"
    prs = Presentation(pptx_file)

    # Get source and destination shapes by placeholder name
    src_shape = find_placeholders(prs.slides[0], "bodytext")[0]
    dst_shape = find_placeholders(prs.slides[1], "bodytext")[0]

    # Apply paragraph-level style
    p = src_shape.text_frame.paragraphs[0]
    set_paragraph_style(
        paragraph=p,
        font_name="Arial",
        font_size=16,
        bold=True,
        italic=False,
        font_color=(0, 102, 204),
        align="center",
    )
    logger.info("Extracted paragraph style: %s", extract_paragraph_style(p))
    prs.save("data/text_style1.pptx")

    # Apply text style to entire text box
    set_textbox_style(
        src_shape,
        font_name="微软雅黑",
        font_size=24,
        font_color=(80, 80, 80),
        bold=False,
        align="left",
    )
    logger.info("Extracted textbox style: %s", extract_textbox_style(src_shape))
    prs.save("data/text_style2.pptx")

    transfer_text_style(src_shape, dst_shape)
    prs.save("data/text_styled_transfer_full.pptx")

    transfer_text_style(src_shape, dst_shape, mode="single")
    prs.save("data/text_styled_transfer_single.pptx")