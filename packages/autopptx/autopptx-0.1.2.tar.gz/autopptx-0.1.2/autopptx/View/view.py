import logging
import os
import uuid
from pptx import Presentation

from autopptx.Type.type import (
    is_text, is_title, is_subtitle, is_bodytext,
    is_image, is_table
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def get_table(shape):
    """Extract all rows of text from a PPT table shape."""
    if not is_table(shape):
        logger.warning(f"⛔️ Shape at index unknown is not a table shape: {getattr(shape, 'shape_type', 'Unknown')}")
        return None

    table_data = []
    for row in shape.table.rows:
        row_data = [cell.text.strip() for cell in row.cells]
        table_data.append(row_data)
    return table_data


def get_text(shape):
    """Extract text content from a text-containing shape."""
    if not is_text(shape):
        logger.warning(f"⛔️ Shape is not a text-containing shape: {shape.shape_type}")
        return None
    return shape.text.strip()


def get_image(shape, output_dir="data/extracted_images"):
    """Extract image from shape and save to local file."""
    if not is_image(shape):
        logger.warning(f"⛔️ Shape is not an image shape: {shape.shape_type}")
        return None

    os.makedirs(output_dir, exist_ok=True)
    image = shape.image
    ext = image.ext or "png"
    filename = f"{uuid.uuid4().hex}.{ext}"
    filepath = os.path.join(output_dir, filename)

    with open(filepath, "wb") as f:
        f.write(image.blob)

    return filepath


def view_slide(slide):
    """
    Extract all content (text, table, image) from a slide.
    Returns a list of dicts with index, type, content.
    """
    results = []
    for idx, shape in enumerate(slide.shapes):
        shape_info = {
            "index": idx,
            "type": None,
            "content": None,
        }

        if is_table(shape):
            shape_info["type"] = "table"
            shape_info["content"] = get_table(shape)

        elif is_image(shape):
            shape_info["type"] = "image"
            shape_info["content"] = get_image(shape)

        elif is_text(shape):
            if is_title(shape):
                shape_info["type"] = "title"
            elif is_subtitle(shape):
                shape_info["type"] = "subtitle"
            elif is_bodytext(shape):
                shape_info["type"] = "bodytext"
            else:
                shape_info["type"] = "text"
            shape_info["content"] = get_text(shape)

        else:
            shape_info["type"] = str(shape.shape_type)
            logger.info(f"ℹ️ Unsupported shape type at index {idx}: {shape.shape_type}")

        results.append(shape_info)

    return results


if __name__ == '__main__':
    ppt_path = "data/output_demo.pptx"
    prs = Presentation(ppt_path)

    for i, slide in enumerate(prs.slides):
        logger.info(f"📄 Slide {i + 1}")
        contents = view_slide(slide)

        for item in contents:
            type_tag = item['type']
            logger.info(f"  ▶ Shape {item['index']} ({type_tag}):")
            if type_tag == "table":
                for row in item["content"]:
                    logger.info(f"     {row}")
            else:
                logger.info(f"     {item['content']}")