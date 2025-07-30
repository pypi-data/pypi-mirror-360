import json
import argparse
from pptx import Presentation

from autopptx.Text.texts import replace_title, replace_subtitle, replace_bodytexts
from autopptx.Image.images import replace_images
from autopptx.Table.tables import replace_tables
from autopptx.Type.find import find_placeholders
from autopptx.Table.style import set_table_style


def load_input_data(json_path):
    """
    Load input data from a JSON file.

    Expected format:
    [
        {
            "title": "Title 1",
            "subtitle": "Subtitle 1",
            "bodytext": ["Paragraph 1", "Paragraph 2"],
            "image": ["path/to/img1.png", ...],
            "table": [[["Header1", "Header2"], ["Row1", "Row2"]]]
        },
        ...
    ]
    """
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return data


def process_presentation(template_path, input_data, output_path):
    """
    Replace all placeholders in the presentation using the input data.

    Parameters:
        template_path (str): Path to the PPTX template file.
        input_data (list): List of dictionaries, one per slide.
        output_path (str): Output file path to save the result.
    """
    prs = Presentation(template_path)

    for idx, slide in enumerate(prs.slides):
        if idx >= len(input_data):
            print(f"⚠️ Slide {idx + 1} skipped: no input data.")
            continue

        data = input_data[idx]
        replace_title(slide, data.get("title", ""))
        replace_subtitle(slide, data.get("subtitle", ""))
        replace_bodytexts(slide, data.get("bodytext", ""))
        replace_images(slide, data.get("image", ""))
        replace_tables(slide, data.get("table", ""))

        table_shapes = find_placeholders(slide, 'table')
        if table_shapes:
            for shape in table_shapes:
                set_table_style(shape, font_name="等线")


    prs.save(output_path)
    print(f"✅ Saved generated PPT file: {output_path}")


def main():
    parser = argparse.ArgumentParser(
        description="AutoPPTX: Automated placeholder replacement for PPTX templates"
    )
    parser.add_argument(
        "--template",
        type=str,
        default="./data/template.pptx",
        help="Path to the PPTX template file",
    )
    parser.add_argument(
        "--input",
        type=str,
        default="./data/input_data.json",
        help="Path to the input JSON file",
    )
    parser.add_argument(
        "--output",
        type=str,
        default="./data/output_demo.pptx",
        help="Path to save the generated PPTX file",
    )

    args = parser.parse_args()
    input_data = load_input_data(args.input)
    process_presentation(args.template, input_data, args.output)


if __name__ == "__main__":
    main()