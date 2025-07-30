import logging
from pptx.util import Inches

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="[%(levelname)s] %(message)s",
)
logger = logging.getLogger(__name__)


def set_image_style(
    shape,
    left=1.0,
    top=1.0,
    width=2.0,
    height=2.0,
    rotation=0.0,
):
    """
    Apply position, size, and rotation to an image shape.

    Parameters:
        shape (pptx.shapes.picture.Picture): The image shape to modify.
        left (float): Horizontal position in inches. Default is 1.0.
        top (float): Vertical position in inches. Default is 1.0.
        width (float): Width in inches. Default is 2.0.
        height (float): Height in inches. Default is 2.0.
        rotation (float): Rotation angle in degrees. Default is 0.0.

    Raises:
        ValueError: If the shape is not a Picture object.
    """
    if not hasattr(shape, "image") or not hasattr(shape, "width"):
        raise ValueError("Shape must be a valid image or image placeholder with content.")

    try:
        shape.left = Inches(left)
        shape.top = Inches(top)
        shape.width = Inches(width)
        shape.height = Inches(height)
        shape.rotation = rotation

        logger.info("✅ Applied image style successfully.")
    except Exception as e:
        logger.error("❌ Failed to apply image style: %s", e)
        raise


def extract_image_style(shape):
    """
    Extract style attributes from a picture shape.

    Parameters:
        shape (pptx.shapes.picture.Picture): The image shape.

    Returns:
        dict: A dictionary containing:
            - left (float): Left position in inches
            - top (float): Top position in inches
            - width (float): Width in inches
            - height (float): Height in inches
            - rotation (float): Rotation angle in degrees
    """
    if not hasattr(shape, "image") or not hasattr(shape, "width"):
        raise ValueError("Shape must be a valid image or image placeholder with content.")

    try:
        style = {
            "left": round(shape.left.inches, 3),
            "top": round(shape.top.inches, 3),
            "width": round(shape.width.inches, 3),
            "height": round(shape.height.inches, 3),
            "rotation": shape.rotation,
        }
        logger.info("✅ Extracted image style successfully.")
        return style
    except Exception as e:
        logger.error("❌ Failed to extract image style: %s", e)
        raise


def transfer_image_style(source_shape, target_shape):
    """
    Transfer style attributes from one picture shape to another.

    Parameters:
        source_shape (pptx.shapes.picture.Picture): The shape to copy style from.
        target_shape (pptx.shapes.picture.Picture): The shape to apply style to.

    Raises:
        ValueError: If either shape is not a Picture object.
    """
    if not hasattr(source_shape, "image") or not hasattr(target_shape, "image"):
        raise ValueError("Both shapes must be image shapes (either Picture or PicturePlaceholder with an image).")

    try:
        style = extract_image_style(source_shape)
        set_image_style(
            target_shape,
            left=style["left"],
            top=style["top"],
            width=style["width"],
            height=style["height"],
            rotation=style["rotation"],
        )
        logger.info("✅ Transferred image style successfully.")
    except Exception as e:
        logger.error("❌ Failed to transfer image style: %s", e)
        raise


if __name__ == "__main__":
    from pptx import Presentation
    from autopptx.Type.find import find_placeholders

    prs = Presentation("data/output_demo.pptx")
    slide = prs.slides[3]
    images = find_placeholders(slide, "image")

    if len(images) >= 2:
        transfer_image_style(images[0], images[1])
    else:
        logger.warning("⚠️ Not enough pictures found on the slide.")

    prs.save("data/output_image_style.pptx")