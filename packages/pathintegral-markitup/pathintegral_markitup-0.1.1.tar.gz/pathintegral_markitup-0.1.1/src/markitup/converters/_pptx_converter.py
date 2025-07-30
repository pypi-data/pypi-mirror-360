import sys
import base64
import os
import io
import re
import html

from typing import BinaryIO, Any, List, Dict
from operator import attrgetter

from ._html_converter import HtmlConverter
from .._base_converter import DocumentConverter, DocumentConverterResult
from .._schemas import StreamInfo, Config, MarkdownChunk
import pptx


ACCEPTED_MAGIC_TYPE_PREFIXES = [
    "application/vnd.openxmlformats-officedocument.presentationml",
]

ACCEPTED_FILE_CATEGORY = [".pptx"]


class PptxConverter(DocumentConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def __init__(self, config: Config):
        super().__init__(config=config)
        self._html_converter = HtmlConverter(config=config)

    def convert(
        self,
        file_stream: BinaryIO,
        stream_info: StreamInfo,
        **kwargs: Any,  # Options to pass to the converter
    ) -> DocumentConverterResult:

        # Perform the conversion
        presentation = pptx.Presentation(file_stream)
        md_content = ""
        chunk_list = []
        slide_num = 0

        for slide in presentation.slides:
            slide_num += 1
            slide_text = ""
            slide_images = []

            def get_shape_content(shape, **kwargs):
                nonlocal slide_text
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""

                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        pass

                    # Prepare the alt, escaping any special characters
                    alt_text = "\n".join([alt_text]) or shape.name
                    alt_text = re.sub(r"[\r\n\[\]]", " ", alt_text)
                    alt_text = re.sub(r"\s+", " ", alt_text).strip()

                    # Create image chunk
                    if 'image' in self.config.modalities:
                        blob = shape.image.blob
                        content_type = shape.image.content_type or "image/png"
                        b64_string = base64.b64encode(blob).decode("utf-8")
                        image_md = f"\n![{alt_text}](data:{content_type};base64,{b64_string})\n"
                        slide_images.append(image_md)
                    else:
                        filename = re.sub(r"\W", "", shape.name) + ".jpg"
                        image_md = f"\n![{alt_text}]({filename})\n"
                        slide_images.append(image_md)

                # Tables
                if self._is_table(shape):
                    slide_text += self._convert_table_to_markdown(shape.table, **kwargs)

                # Charts
                if shape.has_chart:
                    slide_text += self._convert_chart_to_markdown(shape.chart)

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        slide_text += "# " + shape.text.lstrip() + "\n"
                    else:
                        slide_text += shape.text + "\n"

                # Group Shapes
                if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
                    sorted_shapes = sorted(shape.shapes, key=attrgetter("top", "left"))
                    for subshape in sorted_shapes:
                        get_shape_content(subshape, **kwargs)

            title = slide.shapes.title
            sorted_shapes = sorted(slide.shapes, key=attrgetter("top", "left"))
            for shape in sorted_shapes:
                get_shape_content(shape, **kwargs)

            # Add notes if present
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    slide_text += "\n\n### Notes:\n" + notes_frame.text

            # Create text chunk for the slide if there's any text content
            slide_text = slide_text.strip()
            if slide_text:
                chunk_list.append(MarkdownChunk(
                    chunk_modality='text',
                    chunk_id=len(chunk_list),
                    page_id=slide_num - 1,
                    content=slide_text
                ))

            # Create image chunks for the slide
            for image_md in slide_images:
                chunk_list.append(MarkdownChunk(
                    chunk_modality='image',
                    chunk_id=len(chunk_list),
                    page_id=slide_num - 1,
                    content=image_md
                ))

            # Add slide content to full markdown
            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"
            md_content += slide_text
            md_content += "\n".join(slide_images)
            md_content = md_content.strip()

        return DocumentConverterResult(
            markdown=md_content.strip(),
            markdown_chunk_list=chunk_list if self.config.chunk else None,
            config=self.config,
            stream_info=stream_info
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False

    def _convert_table_to_markdown(self, table, **kwargs):
        # Write the table as HTML, then convert it to Markdown
        html_table = "<html><body><table>"
        first_row = True
        for row in table.rows:
            html_table += "<tr>"
            for cell in row.cells:
                if first_row:
                    html_table += "<th>" + html.escape(cell.text) + "</th>"
                else:
                    html_table += "<td>" + html.escape(cell.text) + "</td>"
            html_table += "</tr>"
            first_row = False
        html_table += "</table></body></html>"

        return (
            self._html_converter.convert_string(
                html_table, **kwargs).markdown.strip()
            + "\n"
        )

    def _convert_chart_to_markdown(self, chart):
        try:
            md = "\n\n### Chart"
            if chart.has_title:
                md += f": {chart.chart_title.text_frame.text}"
            md += "\n\n"
            data = []
            category_names = [c.label for c in chart.plots[0].categories]
            series_names = [s.name for s in chart.series]
            data.append(["Category"] + series_names)

            for idx, category in enumerate(category_names):
                row = [category]
                for series in chart.series:
                    row.append(series.values[idx])
                data.append(row)

            markdown_table = []
            for row in data:
                markdown_table.append("| " + " | ".join(map(str, row)) + " |")
            header = markdown_table[0]
            separator = "|" + "|".join(["---"] * len(data[0])) + "|"
            return md + "\n".join([header, separator] + markdown_table[1:])
        except ValueError as e:
            # Handle the specific error for unsupported chart types
            if "unsupported plot type" in str(e):
                return "\n\n[unsupported chart]\n\n"
        except Exception:
            # Catch any other exceptions that might occur
            return "\n\n[unsupported chart]\n\n"
