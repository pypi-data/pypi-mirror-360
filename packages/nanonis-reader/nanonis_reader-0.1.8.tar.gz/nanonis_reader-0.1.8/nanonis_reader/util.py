from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np
import os
import glob
import nanonis_reader as nr
from mpl_toolkits.axes_grid1 import make_axes_locatable
import io
from nanonis_reader.find_value import nearest


class NanonisData:
    def __init__(self, base_path, file_number=None, keyword=None):
        '''
        base_bath: filepath
        fine_number: file number (ex: 0015.sxm -> 15)
        keyword: file name filter (ex: Au111_xxx_0013.sxm -> 'Au')
        '''
        # 파일 확장자별 함수 매핑
        file_handlers = {
            '.sxm': nr.nanonis_sxm.Load,
            '.dat': nr.nanonis_dat.Load,
            '.3ds': nr.nanonis_3ds.Load
        }
        
        # 파일 번호가 주어진 경우
        if isinstance(file_number, (int, str)):
            # 숫자를 4자리 문자열로 변환 (예: 16 -> '0016')
            number_str = str(file_number).zfill(4)
            
            # 키워드가 주어진 경우
            if keyword:
                pattern = os.path.join(base_path, f'*{keyword}*_{number_str}.*')
            else:
                pattern = os.path.join(base_path, f'*_{number_str}.*')
                
            matching_files = glob.glob(pattern)
            
            if not matching_files:
                if keyword:
                    raise ValueError(f"No file found with number {number_str} and keyword '{keyword}'")
                else:
                    raise ValueError(f"No file found with number {number_str}")
            
            if len(matching_files) > 1:
                print(f"Warning: Multiple files found. Files found:")
                for f in matching_files:
                    print(f"- {os.path.basename(f)}")
                print(f"Using the first one: {os.path.basename(matching_files[0])}")
            
            filepath = matching_files[0]
        else:
            filepath = base_path
            
        # 파일 확장자 추출
        _, extension = os.path.splitext(filepath)
        
        # 해당하는 함수 찾아서 실행
        if extension in file_handlers:
            data = file_handlers[extension](filepath)
            # data의 모든 속성을 현재 객체(self)에 복사
            for attr_name in dir(data):
                if not attr_name.startswith('_'):  # private 속성은 제외
                    setattr(self, attr_name, getattr(data, attr_name))
        else:
            raise ValueError(f"Unsupported file extension: {extension}")


class DataToPPT:
    def __init__(self, base_path, keyword=None, output_filename='output.pptx'):
        '''
        base_path: 파일들이 있는 경로
        keyword: 파일 이름 필터 (예: 'Au', 'SiCG')
        output_filename: 생성될 PPT 파일 이름
        '''
        self.base_path = base_path
        self.keyword = keyword
        self.output_filename = output_filename
        self.prs = Presentation()

    def get_scan_parameters(self, data):
        '''
        파일 타입에 따라 적절한 파라미터 추출 함수 호출
        '''
        if data.fname.endswith('.sxm'):
            return self.get_sxm_parameters(data)
        elif data.fname.endswith('.dat'):
            return self.get_dat_parameters(data)
        elif data.fname.endswith('.3ds'):
            return self.get_3ds_parameters(data)
        else:
            raise ValueError(f"Unsupported file type: {data.fname}")

    def get_sxm_parameters(self, data):
        '''
        header에서 자주 사용되는 파라미터들을 추출하는 함수
        반환값: pixels, scan_range, scan_dir, bias, current 등
        '''
        def format_date(date_str):
            '''
            일.월.년 형식을 년.월.일 형식으로 변환
            '''
            try:
                day, month, year = date_str.split('.')
                return f"{year}.{month}.{day}"
            except:
                return date_str  # 파싱 실패시 원본 반환
            
        params = {
            'pixels': data.header['scan_pixels'],
            'range': data.header['scan_range'],
            'direction': data.header['scan_dir'],
            'angle': data.header['scan_angle'],
            'bias': data.header['bias>bias (v)'],
            'current': data.header['z-controller>setpoint'],
            'scan_time': data.header['rec_time'],
            'scan_date': format_date(data.header['rec_date']),
        }
        params['aspect_ratio'] = (params['pixels'][0]/params['pixels'][1])*(params['range'][1]/params['range'][0])
        params['fname'] = data.fname

        return params
        
    def get_dat_parameters(self, data):
        '''
        .dat 파일의 파라미터 추출
        Z rel (m) 포함 여부에 따라 다른 파라미터 반환
        '''
        def format_date(date_str):
            '''
            일.월.년 시:분:초 형식을 년.월.일 시:분:초 형식으로 변환
            '''
            try:
                # 날짜와 시간 분리
                date_part, time_part = date_str.split(' ')
                
                # 날짜 부분 변환
                day, month, year = date_part.split('.')
                formatted_date = f"{year}.{month}.{day}"
                
                # 날짜와 시간 다시 합치기
                return f"{formatted_date}_{time_part}"
            except:
                return date_str  # 파싱 실패시 원본 반환

        if 'Z rel (m)' in data.signals.keys():
            params = {
                'bias': data.header['Bias>Bias (V)'],
                'current': data.header['Z-Controller>Setpoint'],
                # 'sweep_num': data.header['Bias Spectroscopy>Number of sweeps'],
                # 'sweep_num': data.header['Z Spectroscopy>Number of sweeps'],
                'sweep_num': (
                                data.header.get('Z Spectroscopy>Number of sweeps') or
                                data.header.get('Bias Spectroscopy>Number of sweeps') or
                                ''
                            ),
                'offset': data.header['Z Spectroscopy>Initial Z-offset (m)'],
                'sweep_z': data.header['Z Spectroscopy>Sweep distance (m)'],
                # 'comment': data.header['Comment01'],
                'comment': (
                    data.header.get('Comment01') or
                    data.header.get('comment') or
                    data.header.get('Comment') or
                    ''
                ),
                'saved_date': format_date(data.header['Saved Date']),
            }
            # return params
        
        elif 'Frequency (Hz)' in data.signals.keys():
            params ={
                'bias': data.header['Bias>Bias (V)'],
                'current': data.header['Z-Controller>Setpoint'],
                'feedback': data.header['Z-Controller>Controller status'],
                'saved_date': format_date(data.header['Saved Date']),
                'comment': (
                    data.header.get('Comment01') or
                    data.header.get('comment') or
                    data.header.get('Comment') or
                    ''
                ),
            }

        elif data.header['Experiment'] == 'History Data':
            params = {
                'history': data.header['Experiment'],
                'comment': (
                    data.header.get('Comment01') or
                    data.header.get('comment') or
                    data.header.get('Comment') or
                    ''
                ),
                'saved_date': format_date(data.header['Saved Date']),
            }

        elif data.header['Experiment'] == 'LongTerm Data':
            params = {
                'long term chart': data.header['Experiment'],
                'comment': (
                    data.header.get('Comment01') or
                    data.header.get('comment') or
                    data.header.get('Comment') or
                    ''
                ),
                'saved_date': format_date(data.header['Saved Date']),
            }

        else:
            params = {
                'bias': data.header['Bias>Bias (V)'],
                'current': (
                            data.header.get('Z-Controller>Setpoint') or 
                            data.header.get('Current>Current (A)') or 
                            ''
                ),
                'sweep_start': (
                    data.header.get('Bias Spectroscopy>Sweep Start (V)') or
                    ''
                ),
                'sweep_end': (
                    data.header.get('Bias Spectroscopy>Sweep End (V)') or
                    ''
                ),
                # 'sweep_num': data.header['Bias Spectroscopy>Number of sweeps'],
                'sweep_num': (
                    data.header.get('Bias Spectroscopy>Number of sweeps') or
                    ''
                ),
                'comment': (
                    data.header.get('Comment01') or
                    data.header.get('comment') or
                    data.header.get('Comment') or
                    ''
                ),
                'saved_date': format_date(data.header['Saved Date']),
            }
        params['fname'] = data.fname

        return params
        
        

    def get_3ds_parameters(self, data):
        '''
        .3ds 파일의 파라미터 추출
        '''
        def format_date(date_str):
            '''
            일.월.년 시:분:초 형식을 년.월.일 시:분:초 형식으로 변환
            '''
            try:
                # 날짜와 시간 분리
                date_part, time_part = date_str.split(' ')
                
                # 날짜 부분 변환
                day, month, year = date_part.split('.')
                formatted_date = f"{year}.{month}.{day}"
                
                # 날짜와 시간 다시 합치기
                return f"{formatted_date}_{time_part}"
            except:
                return date_str  # 파싱 실패시 원본 반환
            
        # params = {
        #         # 'bias': data.header['Bias>Bias (V)'],
        #         # 'current': data.header['Z-Controller>Setpoint'],
        #         'bias': (data.header.get('Bias>Bias (V)') or ''),
        #         'current': (data.header.get('Z-Controller>Setpoint') or ''),
        #         'pixel': data.header['dim_px'],
        #         'range': data.header['size_xy'],
        #         'comment': data.header['comment'],  
        # }
        # params['fname'] = data.fname
        # params = {
        #     # 'pixels': data.header['scan_pixels'],
        #     # 'range': data.header['scan_range'],
        #     # 'direction': data.header['scan_dir'],
        #     # 'angle': data.header['scan_angle'],
        #     # 'bias': data.header['bias>bias (v)'],
        #     # 'current': data.header['z-controller>setpoint'],
        #     'scan_time': data.header['rec_time'],
        #     'scan_date': format_date(data.header['rec_date']),
        # }
        # params['aspect_ratio'] = (params['pixels'][0]/params['pixels'][1])*(params['range'][1]/params['range'][0])
        # params['fname'] = data.fname

        if all(key not in data.signals.keys() for key in ['LI Demod 1 X (A)', 'LI Demod 1 X [AVG] (A)', 'LI Demod 2 X (A)', 'LI Demod 2 X [AVG] (A)']):
            params = {
                'pixels': data.header['dim_px'],
                'range': data.header['size_xy'],
                'angle': data.header['angle'],
                'bias': (data.header.get('Bias>Bias (V)') or ''),
                'current': (data.header.get('Z-Controller>Setpoint') or ''),
                'sweep_num': (
                                data.header.get('Z Spectroscopy>Number of sweeps') or
                                ''
                            ),
                'offset': (data.header.get('Z Spectroscopy>Initial Z-offset (m)') or ''),
                'sweep_z': (data.header.get('Z Spectroscopy>Sweep distance (m)') or ''),
                'comment': (
                    data.header.get('Comment01') or
                    data.header.get('comment') or
                    data.header.get('Comment') or
                    ''
                ),
                'saved_date': format_date(data.header['start_time']),
            }

        else:
            params = {
                'pixels': data.header['dim_px'],
                'range': data.header['size_xy'],
                'angle': data.header['angle'],
                'bias': (data.header.get('Bias>Bias (V)') or ''),
                'current': (data.header.get('Z-Controller>Setpoint') or ''),
                'sweep_start': (
                    data.header.get('Bias Spectroscopy>Sweep Start (V)') or
                    ''
                ),
                'sweep_end': (
                    data.header.get('Bias Spectroscopy>Sweep End (V)') or
                    ''
                ),
                'sweep_num': (
                    data.header.get('Bias Spectroscopy>Number of sweeps') or
                    ''
                ),
                'comment': (
                    data.header.get('Comment01') or
                    data.header.get('comment') or
                    data.header.get('Comment') or
                    ''
                ),
                'saved_date': format_date(data.header['start_time']),
            }
        params['aspect_ratio'] = (params['pixels'][0]/params['pixels'][1])*(params['range'][1]/params['range'][0])
        params['fname'] = data.fname
        
        return params
    
    def get_sxm_info_text(self, params):
        '''
        .sxm 파일의 정보 텍스트 생성
        '''
        info_texts = []
        info_texts.append(f"{params['fname']}\n")
        info_texts.append(f"{float(params['bias'])} V /")
        current = float(params['current'])
        if abs(current) >= 1e-9:
            # nA 단위로 표시
            info_texts.append(f"{current*1e9:.0f} nA")
        else:
            # pA 단위로 표시
            info_texts.append(f"{current*1e12:.0f} pA")
        info_texts.append(f"\n{params['range'][0]*1e9:.0f} x {params['range'][1]*1e9:.0f} nm²")
        info_texts.append(f"({params['direction']}, {float(params['angle']):.1f}˚)")
        info_texts.append(f"\n({params['scan_date']}_{params['scan_time']})")
        
        return " ".join(info_texts)

    def get_dat_info_text(self, params):
        '''
        .dat 파일의 정보 텍스트 생성
        '''
        if 'sweep_z' in params:
            info_texts = []
            info_texts.append(f"{params['fname']}\n")
            info_texts.append(f"{float(params['bias'])} V /")
            current = float(params['current'])
            if abs(current) >= 1e-9:
                # nA 단위로 표시
                info_texts.append(f"{current*1e9:.0f} nA")
            else:
                # pA 단위로 표시
                info_texts.append(f"{current*1e12:.0f} pA")
            info_texts.append(f"\nComment: {params['comment']}")    
            info_texts.append(f"\n({params['saved_date']})")
        elif 'feedback' in params:
            info_texts = []
            info_texts.append(f"{params['fname']}\n")
            info_texts.append(f"Comment: {params['comment']}")
            info_texts.append(f"\n({params['saved_date']})")
        elif 'history' in params:
            info_texts = []
            info_texts.append(f"{params['fname']}\n")
            info_texts.append(f"{params['history']}")
            info_texts.append(f"\nComment: {params['comment']}")
            info_texts.append(f"\n({params['saved_date']})")
        elif 'long term chart' in params:
            info_texts = []
            info_texts.append(f"{params['fname']}\n")
            info_texts.append(f"{params['long term chart']}")
            info_texts.append(f"\nComment: {params['comment']}")
            info_texts.append(f"\n({params['saved_date']})")
        else:
            info_texts = []
            info_texts.append(f"{params['fname']}\n")
            info_texts.append(f"{float(params['bias'])} V /")
            current = float(params['current'])
            if abs(current) >= 1e-9:
                # nA 단위로 표시
                info_texts.append(f"{current*1e9:.0f} nA")
            else:
                # pA 단위로 표시
                info_texts.append(f"{current*1e12:.0f} pA")
            info_texts.append(f"\n{float(params['sweep_start'])} V to {float(params['sweep_end'])} V (sweeps: {params['sweep_num']})")
            info_texts.append(f"\nComment: {params['comment']}")
            info_texts.append(f"\n({params['saved_date']})")
        
        return " ".join(info_texts)

    def get_3ds_info_text(self, params):
        '''
        .3ds 파일의 정보 텍스트 생성
        '''
        # .3ds 파일에 맞는 정보 포맷
        # For I-z grid,
        if 'sweep_z' in params:
            info_texts = []
            info_texts.append(f"{params['fname']}\n") # Data name
            info_texts.append(f"I-z spectroscopy grid\n") # "I-z spectrum"
            info_texts.append(f"{float(params['bias'])} V /" if params.get('bias') else 'Set bias was not saved.\n') # set bias
            if params.get('current'):
                current = float(params['current'])
                if abs(current) >= 1e-9: # set current
                    # nA 단위로 표시
                    info_texts.append(f"{current*1e9:.0f} nA")
                else:
                    # pA 단위로 표시
                    info_texts.append(f"{current*1e12:.0f} pA")
            else:
                info_texts.append('Set point current was not saved.\n')
            # size (angle)
            # sweep $\Delta z$, number of sweep
            if params.get('offset') and params.get('sweep_z'):
                info_texts.append(f"\n{float(params['offset'])*1e12:.0f} pm to {(float(params['offset'])+float(params['sweep_z']))*1e12:.0f} pm (sweeps: {params['sweep_num']})")
            else:
                info_texts.append('Z sweep range was not saved.\n')
            info_texts.append(f"\n{params['range'][0]*1e9:.0f} x {params['range'][1]*1e9:.0f} nm²")
            info_texts.append(f"({float(params['angle']):.1f}˚)")
            info_texts.append(f"\nComment: {params['comment']}")    
            info_texts.append(f"\n({params['saved_date']})")

        # For STS grid,
        else:
            info_texts = []
            info_texts.append(f"{params['fname']}\n")
            info_texts.append(f"STS grid\n") # "STS"
            info_texts.append(f"{float(params['bias'])} V /")
            current = float(params['current'])
            if abs(current) >= 1e-9:
                # nA 단위로 표시
                info_texts.append(f"{current*1e9:.0f} nA")
            else:
                # pA 단위로 표시
                info_texts.append(f"{current*1e12:.0f} pA")
            info_texts.append(f"\n{float(params['sweep_start'])} V to {float(params['sweep_end'])} V (sweeps: {params['sweep_num']})")
            info_texts.append(f"\n{params['range'][0]*1e9:.0f} x {params['range'][1]*1e9:.0f} nm²")
            info_texts.append(f"({float(params['angle']):.1f}˚)")
            info_texts.append(f"\nComment: {params['comment']}")
            info_texts.append(f"\n({params['saved_date']})")

        # return "3DS file parameters"
        return " ".join(info_texts)

    def get_3sigma_limits(self, data):
        mean = np.nanmean(data)
        sigma = np.nanstd(data)
        return mean + np.array([-3, 3]) * sigma

    def process_sxm_file(self, data):
        '''
        .sxm 파일 처리 함수
        '''
        params = self.get_scan_parameters(data)

        base_size = 5
        figsize = (base_size, base_size)

        # 첫 번째 이미지 (topography)
        fig = plt.figure(figsize=figsize)
        ax = fig.add_subplot(111)      
        
        topo = nr.nanonis_sxm.topography(data)
        z_data = topo.get_z('subtract linear fit', 'fwd')

        origin = 'upper' if params['direction'] == 'down' else 'lower'
        vmin, vmax = self.get_3sigma_limits(z_data)
        nanox = nr.cmap_custom.nanox()
        bwr = nr.cmap_custom.bwr()
        
        # 이미지 플롯
        im = ax.imshow(z_data, origin=origin, vmin=vmin, vmax=vmax, 
                    aspect=params['aspect_ratio'], cmap=nanox, interpolation='none')

        # colorbar 추가
        plt.draw()
        posn = ax.get_position()
        cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 
                            0.02, posn.height])
        plt.colorbar(im, cax=cax)

        # figure를 이미지로 저장
        img_stream1 = io.BytesIO()
        plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.01)
        img_stream1.seek(0)
        # plt.close('all')

        # 두 번째 이미지 (differentiated)
        fig = plt.figure(figsize=figsize)
        ax = fig.add_subplot(111)
        
        z_data_diff = topo.get_z('differentiate', 'fwd')
        vmin, vmax = self.get_3sigma_limits(z_data_diff)
        im = ax.imshow(z_data_diff, origin=origin, vmin=vmin, vmax=vmax, 
                    aspect=params['aspect_ratio'], cmap=nanox, interpolation='none')

        # plt.draw()
        # posn = ax.get_position()
        # cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
        # plt.colorbar(im, cax=cax)

        plt.draw()
        posn = ax.get_position()
        cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
        cbar = plt.colorbar(im, cax=cax)
        cbar.formatter.set_powerlimits((-3, 4))  # scientific notation 사용 범위 설정
        cbar.update_ticks()

        img_stream2 = io.BytesIO()
        plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.01)
        img_stream2.seek(0)
        # plt.close('all')

        if 'LI_Demod_1_X' in data.signals.keys():
            fig = plt.figure(figsize=figsize)
            ax = fig.add_subplot(111)

            didv = nr.nanonis_sxm.didvmap(data)
            didv_data = didv.get_map()
            vmin, vmax = self.get_3sigma_limits(didv_data)
            im = ax.imshow(didv_data, origin=origin, vmin=vmin, vmax=vmax, 
                        aspect=params['aspect_ratio'], cmap=bwr, interpolation='none')

            plt.draw()
            posn = ax.get_position()
            cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
            plt.colorbar(im, cax=cax)

            img_stream3 = io.BytesIO()
            plt.savefig(img_stream3, format='png', bbox_inches='tight', pad_inches=0.01)
            img_stream3.seek(0)
            # plt.close('all')
            
            return img_stream1, img_stream2, img_stream3

        return img_stream1, img_stream2
    
    def process_dat_file(self, data):
        '''
        .dat 파일 처리 함수
        '''
        params = self.get_dat_parameters(data)
        base_size = 5
        figsize = (base_size, base_size)

        if 'sweep_z' in params:
            spec = nr.nanonis_dat.z_spectrum(data)

            # Z-I linear
            spec_z = spec.get_iz()
            fig = plt.figure(figsize=figsize)
            plt.plot(spec_z[0] * 1e9, spec_z[1] * 1e9, 'k-')
            plt.xlabel('Z (nm)')
            plt.ylabel('Current (nA)')
            img_stream1 = io.BytesIO()
            plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream1.seek(0)
            # plt.close('all')

            # Z-I log
            fig = plt.figure(figsize=figsize)
            plt.plot(spec_z[0] * 1e9, np.abs(spec_z[1] * 1e9), 'k-')
            plt.xlabel('Z (nm)')
            plt.ylabel('|Current| (nA)')
            plt.yscale('log')
            plt.grid(True)
            img_stream2 = io.BytesIO()
            plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream2.seek(0)
            # plt.close('all')

            return img_stream1, img_stream2
        
        elif 'Frequency (Hz)' in data.signals.keys():
            spec = nr.nanonis_dat.noise_spectrum(data)

            # Noise spectrum
            spec_noise = spec.get_noise()
            if 'Current PSD (A/sqrt(Hz))' in data.signals.keys():
                plt.plot(spec_noise[0], spec_noise[1]*1e15, 'k-')
                plt.xlabel('Frequency (Hz)')
                plt.ylabel('Current (fA)')
                plt.yscale('log')
                plt.grid(True)
            else:
                plt.plot(spec_noise[0], spec_noise[1]*1e12, 'k-')
                plt.xlabel('Frequency (Hz)')
                plt.ylabel('Z (pm)')
                plt.yscale('log')
                plt.grid(True)
            img_stream1 = io.BytesIO()
            plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream1.seek(0)

            return img_stream1

        elif data.header['Experiment'] == 'History Data':
            hist = nr.nanonis_dat.history_data(data)

            # Current history
            hist_I = hist.get_history('Current (A)')
            fig = plt.figure(figsize=figsize)
            plt.plot(hist_I[0] * 1e-3, hist_I[1] * 1e9, 'k-')
            plt.xlabel('Time (s)')
            plt.ylabel('Current (nA)')
            img_stream1 = io.BytesIO()
            plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream1.seek(0)
            # plt.close('all')

            # Height history
            hist_z = hist.get_history('Z (m)')
            fig = plt.figure(figsize=figsize)
            plt.plot(hist_z[0] * 1e-3, hist_z[1] * 1e9, 'k-')
            plt.xlabel('Time (s)')
            plt.ylabel('Z (nm)')
            # plt.yscale('log')
            # plt.grid(True)
            img_stream2 = io.BytesIO()
            plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream2.seek(0)
            # plt.close('all')

            return img_stream1, img_stream2

        elif data.header['Experiment'] == 'LongTerm Data':
            LTchart = nr.nanonis_dat.longterm_data(data)

            t_LTchart, z_LTchart = LTchart.get_z_longterm_chart()
            fig = plt.figure(figsize=figsize)
            plt.plot(t_LTchart, z_LTchart * 1e9, 'k-')
            plt.xlabel('Rel. Time (s)')
            plt.ylabel('Z (nm)')
            img_stream1 = io.BytesIO()
            plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream1.seek(0)
            # plt.close('all')

            return img_stream1

        else:
            spec = nr.nanonis_dat.spectrum(data)
            bias_values = spec.iv_raw()[0]
            signs = np.sign(bias_values)

            if np.any(signs < 0) and np.any(signs > 0):
                # Scaled dI/dV
                didv_scaled = spec.didv_scaled()
                fig = plt.figure(figsize=figsize)
                plt.plot(didv_scaled[0], didv_scaled[1] * 1e9, 'k-')
                plt.xlabel('Bias (V)')
                plt.ylabel('dI/dV (nS)')
                img_stream1 = io.BytesIO()
                plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.1)
                img_stream1.seek(0)
                # plt.close('all')

                # Normalized dI/dV
                didv_norm = spec.didv_normalized()
                fig = plt.figure(figsize=figsize)
                plt.plot(didv_norm[0], didv_norm[1], 'k-')
                plt.xlabel('Bias (V)')
                plt.ylabel('Norm. dI/dV')
                img_stream2 = io.BytesIO()
                plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.1)
                img_stream2.seek(0)
                # plt.close('all')

                # I-V
                fig = plt.figure(figsize=figsize)
                iv = spec.iv_raw()
                plt.plot(iv[0], iv[1] * 1e12, 'k-')
                plt.xlabel('Bias (V)')
                plt.ylabel('Current (pA)')
                img_stream3 = io.BytesIO()
                plt.savefig(img_stream3, format='png', bbox_inches='tight', pad_inches=0.1)
                img_stream3.seek(0)
                # plt.close('all')


                return img_stream1, img_stream2, img_stream3

            else:
                # dI/dV raw
                didv_raw = spec.didv_raw()
                fig = plt.figure(figsize=figsize)
                plt.plot(didv_raw[0], didv_raw[1] * 1e12, 'k-')
                plt.xlabel('Bias (V)')
                plt.ylabel('dI/dV (arb.)')
                img_stream1 = io.BytesIO()
                plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.1)
                img_stream1.seek(0)
                # plt.close('all')

                # I-V
                iv = spec.iv_raw()
                fig = plt.figure(figsize=figsize)
                plt.plot(iv[0], iv[1] * 1e12, 'k-')
                plt.xlabel('Bias (V)')
                plt.ylabel('Current (pA)')
                img_stream2 = io.BytesIO()
                plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.1)
                img_stream2.seek(0)
                # plt.close('all')

                if 'Z (m)' in data.signals.keys():
                    # dZ/dV
                    dzdv = spec.dzdv_numerical()
                    fig = plt.figure(figsize=figsize)
                    plt.plot(dzdv[0], dzdv[1], 'k-')
                    plt.xlabel('Bias (V)')
                    plt.ylabel('dZ/dV (nm/V)')
                    img_stream3 = io.BytesIO()
                    plt.savefig(img_stream3, format='png', bbox_inches='tight', pad_inches=0.1)
                    img_stream3.seek(0)
                    # plt.close('all')

                    return img_stream1, img_stream2, img_stream3
                
                return img_stream1, img_stream2

    
    def process_3ds_file(self, data):
        '''
        .3ds 파일 처리 함수
        '''
        params = self.get_3ds_parameters(data)
        base_size = 5
        figsize = (base_size, base_size)
        
        # For I-z spectra,
        if 'sweep_z' in params:
            # 첫 번째 이미지 (topography)
            fig = plt.figure(figsize=figsize)
            ax = fig.add_subplot(111)      
            topo = nr.nanonis_3ds.Topo(data)
            z_data = topo.get_z ('subtract linear fit')
            origin = 'lower'
            vmin, vmax = self.get_3sigma_limits(z_data)
            nanox = nr.cmap_custom.nanox()
            bwr = nr.cmap_custom.bwr()
            
            # 이미지 플롯
            im = ax.imshow(z_data, origin=origin, vmin=vmin, vmax=vmax, 
                        aspect=params['aspect_ratio'], cmap=nanox, interpolation='none')

            # colorbar 추가
            plt.draw()
            posn = ax.get_position()
            cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 
                                0.02, posn.height])
            plt.colorbar(im, cax=cax)

            # figure를 이미지로 저장
            img_stream1 = io.BytesIO()
            plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.01)
            img_stream1.seek(0)
            # plt.close('all')

            # 두 번째 이미지 (current map)
            fig = plt.figure(figsize=figsize)
            ax = fig.add_subplot(111)
            spec = nr.nanonis_3ds.Map(data)

            idx = (nearest (spec.signals['sweep_signal'], -100e-12)[0])

            spec_z = spec.get_currentmap (sweep_idx=idx)
            vmin, vmax = self.get_3sigma_limits(spec_z)
            im = ax.imshow(spec_z, origin=origin, vmin=vmin, vmax=vmax, 
                        aspect=params['aspect_ratio'], cmap=nanox, interpolation='none')

            # colorbar 추가
            plt.draw()
            posn = ax.get_position()
            cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
            cbar = plt.colorbar(im, cax=cax)
            cbar.formatter.set_powerlimits((-3, 4))  # scientific notation 사용 범위 설정
            cbar.update_ticks()

            # figure를 이미지로 저장
            img_stream2 = io.BytesIO()
            plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.01)
            img_stream2.seek(0)
            # plt.close('all')

            # 세 번째 이미지 (barrier map)
            fig = plt.figure(figsize=figsize)
            ax = fig.add_subplot(111)
            spec = nr.nanonis_3ds.Map(data)
            barrier_height = spec.get_apparent_barrier_height_map ()[0]
            vmin, vmax = self.get_3sigma_limits(barrier_height)
            im = ax.imshow(barrier_height, origin=origin, vmin=vmin, vmax=vmax, 
                        aspect=params['aspect_ratio'], cmap=bwr, interpolation='none')

            # colorbar 추가
            plt.draw()
            posn = ax.get_position()
            cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
            cbar = plt.colorbar(im, cax=cax)
            cbar.formatter.set_powerlimits((-3, 4))  # scientific notation 사용 범위 설정
            cbar.update_ticks()

            # figure를 이미지로 저장
            img_stream3 = io.BytesIO()
            plt.savefig(img_stream3, format='png', bbox_inches='tight', pad_inches=0.01)
            img_stream3.seek(0)
            # plt.close('all')

            return img_stream1, img_stream2, img_stream3

        # For STS grid,
        else:
            # 첫 번째 이미지 (topography)
            fig = plt.figure(figsize=figsize)
            ax = fig.add_subplot(111)      
            topo = nr.nanonis_3ds.Topo(data)
            z_data = topo.get_z ('subtract linear fit')
            origin = 'lower'
            vmin, vmax = self.get_3sigma_limits(z_data)
            nanox = nr.cmap_custom.nanox()
            bwr = nr.cmap_custom.bwr()
            
            # 이미지 플롯
            im = ax.imshow(z_data, origin=origin, vmin=vmin, vmax=vmax, 
                        aspect=params['aspect_ratio'], cmap=nanox, interpolation='none')

            # colorbar 추가
            plt.draw()
            posn = ax.get_position()
            cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 
                                0.02, posn.height])
            plt.colorbar(im, cax=cax)

            # figure를 이미지로 저장
            img_stream1 = io.BytesIO()
            plt.savefig(img_stream1, format='png', bbox_inches='tight', pad_inches=0.01)
            img_stream1.seek(0)
            # plt.close('all')

            # 두 번째 이미지 (dI/dV map)
            fig = plt.figure(figsize=figsize)
            ax = fig.add_subplot(111)
            spec = nr.nanonis_3ds.Map(data)

            idx = (nearest (spec.signals['sweep_signal'], -0.3)[0] or
                   nearest (spec.signals['sweep_signal'], 8.2)[0] )

            didvmap = spec.get_didvmap (sweep_idx=idx)
            vmin, vmax = self.get_3sigma_limits(didvmap)
            im = ax.imshow(didvmap, origin=origin, vmin=vmin, vmax=vmax, 
                        aspect=params['aspect_ratio'], cmap=bwr, interpolation='none')

            # colorbar 추가
            plt.draw()
            posn = ax.get_position()
            cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
            cbar = plt.colorbar(im, cax=cax)
            cbar.formatter.set_powerlimits((-3, 4))  # scientific notation 사용 범위 설정
            cbar.update_ticks()

            # figure를 이미지로 저장
            img_stream2 = io.BytesIO()
            plt.savefig(img_stream2, format='png', bbox_inches='tight', pad_inches=0.01)
            img_stream2.seek(0)
            # plt.close('all')

            # # 세 번째 이미지 (individual and average STS curves)
            # fig = plt.figure(figsize=figsize)
            # ax = fig.add_subplot(111)
            # spec = nr.nanonis_3ds.Map(data)
            # barrier_height = spec.get_apparent_barrier_height_map ('AVG')[0]
            # vmin, vmax = self.get_3sigma_limits(barrier_height)
            # im = ax.imshow(barrier_height, origin=origin, vmin=vmin, vmax=vmax, 
            #             aspect=params['aspect_ratio'], cmap=bwr, interpolation='none')

            # # colorbar 추가
            # plt.draw()
            # posn = ax.get_position()
            # cax = fig.add_axes([posn.x1 + 0.01, posn.y0, 0.02, posn.height])
            # cbar = plt.colorbar(im, cax=cax)
            # cbar.formatter.set_powerlimits((-3, 4))  # scientific notation 사용 범위 설정
            # cbar.update_ticks()

            # # figure를 이미지로 저장
            # img_stream3 = io.BytesIO()
            # plt.savefig(img_stream3, format='png', bbox_inches='tight', pad_inches=0.01)
            # img_stream3.seek(0)
            # # plt.close('all')

            # Average STS
            fig = plt.figure(figsize=figsize)
            ax = fig.add_subplot(111)
            spec = nr.nanonis_3ds.PtSpec(data)
            lines, pixels = np.shape(z_data)

            # 모든 didv 데이터를 한 번에 수집
            # grid를 얻는 도중 멈출 경우, 나머지 data는 [0,0, ..., 0] -> [nan, nan, ..., nan]으로 처리.
            didv_data = [np.where( np.all( (data := spec.get_didv_raw(i, j, channel='LI Demod 1 X (A)')[1]) == 0 ), np.nan, data ) 
                        for i in range(lines) 
                        for j in range(pixels)]

            # 전압 데이터는 한 번만 가져옴 (모든 지점에서 동일함)
            v = spec.get_didv_raw(0, 0, channel='LI Demod 1 X (A)')[0]

            # 모든 곡선을 한 번에 플롯
            plt.plot(v, np.array(didv_data).T * 1e9, 'k-', alpha=0.2, lw=0.2)

            # 평균 곡선 계산 및 플롯
            didv_avg = np.nanmean(didv_data, axis=0)
            plt.plot(v, didv_avg * 1e9, 'r-')

            plt.xlabel('Bias (V)')
            plt.ylabel('dI/dV (a.u.)')
            img_stream3 = io.BytesIO()
            plt.savefig(img_stream3, format='png', bbox_inches='tight', pad_inches=0.1)
            img_stream3.seek(0)
            # plt.close('all')

            return img_stream1, img_stream2, img_stream3
        
        img_stream = io.BytesIO()
        plt.savefig(img_stream, format='png', bbox_inches='tight')
        img_stream.seek(0)
        # plt.close('all')
        return img_stream
    

    def add_slide(self, data):
        '''
        데이터를 처리하고 슬라이드에 추가하는 함수
        '''
        # 새 슬라이드 추가
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # 빈 슬라이드
        
        # 제목 추가
        title_shape = slide.shapes.title
        title_shape.text = f"File: {data.fname}"
        
        # 파일 확장자에 따른 처리
        if data.fname.endswith('.sxm'):
            params = self.get_scan_parameters(data)
            
            base_size = 3.2  # 이미지 크기를 조금 줄여서 3개가 들어갈 수 있게 조정
            width = Inches(base_size)
            height = Inches(base_size)
            img_top = Inches(1.5)
            
            img_streams = self.process_sxm_file(data)
            if not isinstance(img_streams, tuple):
                img_streams = (img_streams,)
            for i, img_stream in enumerate(img_streams):
                x_position = Inches(0 + base_size * 1.02 * i)
                slide.shapes.add_picture(img_stream, x_position, img_top, width=width)

            text_top = img_top + height + Inches(0.2)  # 0.2인치 간격
            info_text = self.get_sxm_info_text(params)

        elif data.fname.endswith('.dat'):
            params = self.get_dat_parameters(data)
            
            base_size = 3.2
            width = Inches(base_size)
            height = Inches(base_size)
            img_top = Inches(1.5)

            img_streams = self.process_dat_file(data)
            if not isinstance(img_streams, tuple):
                img_streams = (img_streams,)
            for i, img_stream in enumerate(img_streams):
                x_position = Inches(0 + base_size * 1.02 * i)
                slide.shapes.add_picture(img_stream, x_position, img_top, width=width)
                

            text_top = img_top + height + Inches(0.2)  # 0.2인치 간격
            info_text = self.get_dat_info_text(params)
            
        elif data.fname.endswith('.3ds'):
            params = self.get_3ds_parameters(data)
            # info_text = self.get_3ds_info_text(params)
            # text_top = Inches(6)
            # img_stream = self.process_3ds_file(data)
            base_size = 3.2
            width = Inches(base_size)
            height = Inches(base_size)
            img_top = Inches(1.5)

            img_streams = self.process_3ds_file(data)
            if not isinstance(img_streams, tuple):
                img_streams = (img_streams,)
            for i, img_stream in enumerate(img_streams):
                x_position = Inches(0 + base_size * 1.02 * i)
                slide.shapes.add_picture(img_stream, x_position, img_top, width=width)
                

            text_top = img_top + height + Inches(0.2)  # 0.2인치 간격
            info_text = self.get_3ds_info_text(params)
        
        # 추가 정보 텍스트 박스
        # left, top, width, height
        txBox = slide.shapes.add_textbox(Inches(1), text_top,
                                    Inches(8), Inches(0.5))
        tf = txBox.text_frame
        tf.text = info_text if info_text else "No parameters available"

    def find_max_file_number(self):
        '''
        현재 경로에서 가장 큰 파일 번호 찾기
        '''
        import glob
        import os

        pattern = os.path.join(self.base_path, f'*{self.keyword if self.keyword else ""}*.*')
        files = glob.glob(pattern)
        
        max_num = 0
        for file in files:
            try:
                # 파일 이름에서 숫자 추출
                num_str = file.split('_')[-1].split('.')[0]
                num = int(num_str)
                max_num = max(max_num, num)
            except:
                continue
                
        return max_num

    def generate_ppt(self):
        '''
        PPT 생성 메인 함수
        '''
        # 최대 파일 번호 찾기
        max_num = self.find_max_file_number()
        print(f"\nMaximum file number in directory: {max_num}")
        
        # 사용자 입력 받기
        while True:
            try:
                start = input("\nEnter start number (or 'q' to quit): ")
                if start.lower() == 'q':
                    print("Cancelled by user")
                    return
                start = int(start)
                
                end = input("Enter end number (or 'q' to quit): ")
                if end.lower() == 'q':
                    print("Cancelled by user")
                    return
                end = int(end)
                
                if start > end:
                    print("Start number should be less than end number")
                    continue
                if end > max_num:
                    print(f"Warning: End number ({end}) is larger than maximum file number ({max_num})")
                    proceed = input("Do you want to proceed anyway? (y/n): ")
                    if proceed.lower() != 'y':
                        continue
                
                break
            except ValueError:
                print("Please enter valid numbers (or 'q' to quit)")
        
        proceed = input(f"\nGenerate PPT for files {start} to {end}? (y/n): ")
        if proceed.lower() != 'y':
            print("Operation cancelled")
            return
            
        print(f"\nGenerating PPT for files {start} to {end}...")
        
        for i in range(start, end + 1):
            try:
                # 파일 로드
                data = NanonisData(self.base_path, i, self.keyword)
                print(f"Processing file: {data.fname}")
                
                # 슬라이드 추가
                self.add_slide(data)
                plt.close('all')
                
            except ValueError as e:
                print(f"Skipping number {i}: {str(e)}")
                plt.close('all')
                continue
        
        # PPT 저장
        save_path = self.base_path + 'PPT/'
        if not os.path.exists(save_path):
            os.makedirs(save_path)
        self.prs.save(save_path + self.output_filename)
        print(f"\nPPT has been saved as {save_path + self.output_filename}")