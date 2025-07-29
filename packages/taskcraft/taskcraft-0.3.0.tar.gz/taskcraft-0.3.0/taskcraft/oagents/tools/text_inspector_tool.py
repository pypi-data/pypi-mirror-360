'''
@ TOOL for Text Inspection
'''

from typing import Optional
from xml.dom import minidom
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
import json
from Bio import PDB
from pptx import Presentation


from .tools import Tool
from ..models import MessageRole, Model
from ..mdconvert import MarkdownConverter


class TextInspectorTool(Tool):
    name = "inspect_file_as_text"
    description = """
You cannot load files yourself: instead call this tool to read a file as markdown text and ask questions about it.
This tool handles the following file extensions: [".html", ".pdb", ".xlsx", ".xls", ".pdf", ".docx", ".ppt", ".pptx"], and all other types of text files. IT DOES NOT HANDLE IMAGES."""

    inputs = {
        "file_path": {
            "description": "The path to the file you want to read as text. Must be a '.something' file, like '.pdf'. If it is an image, use the visualizer tool instead! If it is an audio, use the audio tool instead! DO NOT use this tool for an HTML webpage: use the web_search tool instead!",
            "type": "string",
        },
        "question": {
            "description": "[Optional]: Your question, as a natural language sentence. Provide as much context as possible. Do not pass this parameter if you just want to directly return the content of the file.",
            "type": "string",
            "nullable": True,
        },
    }
    output_type = "string"
    md_converter = MarkdownConverter()

    def __init__(self, model: Model, text_limit: int):
        super().__init__()
        self.model = model
        self.text_limit = text_limit

    def jsonld_to_markdown(self, data):
        """
        将 JSON-LD 数据转换为 Markdown 格式
        """
        markdown = ""
        if isinstance(data, dict):
            for key, value in data.items():
                markdown += f"**{key}**: {self.jsonld_to_markdown(value)}\n"
        elif isinstance(data, list):
            for item in data:
                markdown += f"- {self.jsonld_to_markdown(item)}\n"
        else:
            markdown += str(data)
        return markdown
    
    def parse_pdb_file(self, file_path):
        """
        author: 成冠桥
        TODO: 解析 PDB 文件，提取基本的分子结构信息。
        :param file_path: pdb文件路径
        """
        parser = PDB.PDBParser(QUIET=True)
        structure = parser.get_structure("protein", file_path)

        atoms = list(structure.get_atoms())
        if len(atoms) < 2:
            return "Error: PDB file contains fewer than two atoms."

        atom1, atom2 = atoms[0], atoms[1]
        distance = atom1 - atom2  # 计算原子间的欧几里得距离

        return f"First atom: {atom1.get_name()} ({atom1.coord})\n" \
            f"Second atom: {atom2.get_name()} ({atom2.coord})\n" \
            f"Distance: {distance:.3f} Angstroms ({distance * 100:.0f} pm)"

    def extract_excel_data(self, file_path, max_rows=30):
        try:
            workbook = load_workbook(file_path)
            result = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                for row in sheet.iter_rows():
                    row_data = []
                    for cell in row:
                        cell_value = cell.value if cell.value is not None else ""

                        fill = cell.fill
                        if hasattr(fill, "fgColor") and fill.fgColor.type == "rgb":
                            cell_color = fill.fgColor.rgb
                            if len(cell_color) == 8:
                                cell_color = cell_color[2:]
                        else:
                            cell_color = "FFFFFF"
                        row_data.append({
                            "value": cell_value,
                            "color": cell_color
                        })
                    result.append(row_data)
                json_data = json.dumps(result, ensure_ascii=False, indent=4)
                data_with_colors = json.loads(json_data)
                text = []
                num_rows = len(data_with_colors)
                num_cols = len(data_with_colors[0]) if data_with_colors else 0
                text.append(f"This is a {num_rows} rows and {num_cols} columns table. The content is shown below:")
                for i, row in enumerate(data_with_colors):
                    row_text = ""
                    for j, cell in enumerate(row):
                        if cell["value"] != "":
                            value = cell["value"]
                        else:
                            value = "None"
                        color = cell["color"]
                        if color == "FFFFFF" or color == "000000":
                            row_text += f"{value} "
                        else:
                            row_text += f"{value}({color}) "
                    text.append(row_text)
            return "\n".join(text)
        except Exception as e:
            return f"Error: {str(e)}"

    def forward_initial_exam_mode(self, file_path, question):
        result = self.md_converter.convert(file_path)
        # 添加 XML 解析支持 from成冠桥
        if file_path.endswith(".xml"):
            try:
                dom = minidom.parse(file_path)
                result_text = " ".join(
                    [node.firstChild.nodeValue for node in dom.getElementsByTagName("*") if node.firstChild and node.firstChild.nodeType == node.TEXT_NODE]
                )
                result.text_content = result_text
            except Exception as e:
                raise Exception(f"Error parsing XML file: {str(e)}")

        if file_path.endswith(".csv"):
            try:
                with open(file_path, 'r') as fr:
                    contents = fr.readlines()
                result.text_content = contents
            except Exception as e:
                raise Exception(f"Error parsing CSV file: {str(e)}")

        if file_path.endswith(".pdb"):
            try:
                pdb_info = self.parse_pdb_file(file_path)
                if not question:
                    return f"Extracted PDB Data:\n{pdb_info}"
                # return f"Extracted PDB Data:\n{pdb_info}\n\nQuestion: {question}\nAnswer: (Model-generated answer based on extracted PDB data)"
                else:
                    result.text_content = pdb_info
            except Exception as e:
                raise Exception(f"Error parsing PDB file: {e}")

        if file_path.endswith((".ppt",  ".pptx")):
            content = ""
            try:
                ppt = Presentation(file_path)
                for slide_number, slide in enumerate(ppt.slides, start=1):
                    content += f"=== Slide {slide_number} ===\n"
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    content += '\n'.join(slide_texts) + '\n\n'
                result.text_content = content.strip()
            except Exception as e:
                return f"Error parsing PPT file: {e}"

        if file_path.endswith((".xls", ".xlsx")):
            try:
                full_content = self.extract_excel_data(file_path)
                if not question:
                    return full_content
                else:
                    result.text_content = full_content
            except Exception as e:
                raise Exception(f"Error processing Excel file: {e}")

        if file_path[-4:] in [".png", ".jpg"]:
            raise Exception("Cannot use inspect_file_as_text tool with images: use visualizer instead!")

        if ".zip" in file_path:
            return result.text_content
            
        if not question:
            return result.text_content

        if len(result.text_content) < 4000:
            return "Document content: " + result.text_content

        messages = [
            {
                "role": MessageRole.SYSTEM,
                "content": [
                    {
                        "type": "text",
                        "text": "Here is a file:\n### "
                        + str(result.title)
                        + "\n\n"
                        + result.text_content[: self.text_limit],
                    }
                ],
            },
            {
                "role": MessageRole.USER,
                "content": [
                    {
                        "type": "text",
                        "text": "Now please write a short, 5 sentence caption for this document, that could help someone asking this question: "
                        + question
                        + "\n\nDon't answer the question yourself! Just provide useful notes on the document",
                    }
                ],
            },
        ]
        return self.model(messages).content

    def forward(self, file_path, question: Optional[str] = None) -> str:
        # result = self.md_converter.convert(file_path)

        if ".jsonld" in file_path:
            try:
                with open(file_path, 'r') as f:
                    data = json.load(f)
                result_text = self.jsonld_to_markdown(data)
                result = type('Result', (object,), {'title': file_path, 'text_content': result_text})
            except Exception as e:
                raise Exception(f"Error parsing JSON-LD file: {str(e)}")
        else:
            result = self.md_converter.convert(file_path)

        if file_path[-4:] in [".png", ".jpg"]:
            raise Exception("Cannot use inspect_file_as_text tool with images: use visualizer instead!")

        if ".zip" in file_path:
            return result.text_content

        if file_path.endswith((".ppt",  ".pptx")):
            content = ""
            try:
                ppt = Presentation(file_path)
                for slide_number, slide in enumerate(ppt.slides, start=1):
                    content += f"=== Slide {slide_number} ===\n"
                    slide_texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())
                    content += '\n'.join(slide_texts) + '\n\n'
                result.text_content = content.strip()
            except Exception as e:
                return f"Error parsing PPT file: {e}"

        if file_path.endswith((".xls", ".xlsx")):
            try:
                full_content = self.extract_excel_data(file_path)
                if not question:
                    return full_content
                else:
                    result.text_content = full_content
            except Exception as e:
                raise Exception(f"Error processing Excel file: {e}")

        if file_path.endswith(".pdb"):
            try:
                pdb_info = self.parse_pdb_file(file_path)
                if not question:
                    return f"Extracted PDB Data:\n{pdb_info}"
                # return f"Extracted PDB Data:\n{pdb_info}\n\nQuestion: {question}\nAnswer: (Model-generated answer based on extracted PDB data)"
                else:
                    result.text_content = pdb_info
            except Exception as e:
                raise Exception(f"Error parsing PDB file: {e}")
            
        if not question:
            return result.text_content

        messages = [
            {
                "role": MessageRole.SYSTEM,
                "content": [
                    {
                        "type": "text",
                        "text": "You will have to write a short caption for this file, then answer this question:"
                        + question,
                    }
                ],
            },
            {
                "role": MessageRole.USER,
                "content": [
                    {
                        "type": "text",
                        "text": "Here is the complete file:\n### "
                        + str(result.title)
                        + "\n\n"
                        + result.text_content[: self.text_limit],
                    }
                ],
            },
            {
                "role": MessageRole.USER,
                "content": [
                    {
                        "type": "text",
                        "text": "Now answer the question below. Use these three headings: '1. Short answer', '2. Extremely detailed answer', '3. Additional Context on the document and question asked'."
                        + question,
                    }
                ],
            },
        ]
        return self.model(messages).content