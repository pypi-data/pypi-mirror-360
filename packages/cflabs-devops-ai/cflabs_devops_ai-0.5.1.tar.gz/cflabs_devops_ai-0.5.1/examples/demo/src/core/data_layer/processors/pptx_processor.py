import hashlib
from typing import List, Dict, Any, BinaryIO
from pathlib import Path
from pptx import Presentation
from src.base_classes.data_layer_class import DocumentProcessor
from src.core.rag_utils import chunk_document
from src.utils.logger import get_logger, time_function, log_function_call

logger = get_logger(__name__)

class PptxProcessor(DocumentProcessor):
    """Processor for PPTX presentations"""
    
    def can_process(self, file_extension: str) -> bool:
        return file_extension.lower() in ['.pptx', '.ppt']
    
    def get_supported_extensions(self) -> List[str]:
        return ['.pptx', '.ppt']
    
    @time_function
    @log_function_call
    def process(self, file: BinaryIO, filename: str) -> List[Dict[str, Any]]:
        """Process PPTX file and return structured document data"""
        logger.info(f"Processing PPTX file: {filename}")
        
        try:
            # Load the presentation
            prs = Presentation(file)
            
            # Extract text from slides
            content_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = [f"Slide {slide_num}:"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # More than just the slide number
                    content_parts.append("\n".join(slide_text))
            
            content = "\n\n".join(content_parts)
            
            if not content.strip():
                logger.warning(f"Empty PPTX file: {filename}")
                return []
            
            # Generate original_uuid as hash of content
            content_hash = hashlib.sha256(content.encode()).hexdigest()
            
            # Create document
            formatted_doc = {
                "doc_id": f"{filename}_doc_1",
                "original_uuid": content_hash,
                "content": content,
                "source": filename,
                "file_type": "pptx"
            }
            
            # Add chunks to the document
            chunked_doc = chunk_document(formatted_doc)
            
            logger.info(f"PPTX processing completed. Documents: 1, Total chunks: {len(chunked_doc['chunks'])}")
            return [chunked_doc]
            
        except Exception as e:
            logger.error(f"Error processing PPTX file {filename}: {str(e)}")
            raise 