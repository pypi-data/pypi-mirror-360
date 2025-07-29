import email
import os
import re
import unicodedata
import zipfile

import docx
import fitz
import numpy as np
import openpyxl
import torch
from bs4 import BeautifulSoup
from googletrans import Translator
from langdetect import detect
from mmh3 import hash
from pptx import Presentation

from ncrypt.utils import (
    EMAIL_EXTENSIONS,
    MAC_EXTENSIONS,
    MARKUP_EXTENSIONS,
    OFFICE_DOCX_EXTENSIONS,
    OFFICE_PPTX_EXTENSIONS,
    OFFICE_XLSX_EXTENSIONS,
    PDF_EXTENSIONS,
    PLAINTEXT_EXTENSIONS,
    RTF_EXTENSIONS,
    TEXT_EMBED_MODEL,
    TEXT_EXTENSIONS,
    TEXT_SUMMARY_MODEL,
    ProcessingError,
    UnsupportedExtensionError,
)

from .image import quantize_embedding

os.environ["TOKENIZERS_PARALLELISM"] = "false"


def _estimate_token_count(text: str) -> int:
    """
    Estimate token count (roughly 1 token per 0.75 words in English).
    """
    words = text.split()

    return int(len(words) / 0.75)


def extract_raw_text(path: str, extension: str) -> str:
    if extension not in TEXT_EXTENSIONS:
        raise UnsupportedExtensionError(f"The provided file extension is not supported for metadata: .{extension}")

    if extension in PLAINTEXT_EXTENSIONS or RTF_EXTENSIONS:
        with open(path, encoding="utf-8", errors="ignore") as file:
            return file.read()

    elif extension in MARKUP_EXTENSIONS:
        with open(path, encoding="utf-8", errors="ignore") as file:
            soup = BeautifulSoup(file, "html.parser")

            return soup.get_text()

    elif extension in PDF_EXTENSIONS:
        with fitz.open(path) as file:
            return "\n".join(page.get_text() for page in file)

    elif extension in OFFICE_DOCX_EXTENSIONS:
        doc = docx.Document(path)

        return "\n".join(p.text for p in doc.paragraphs)

    elif extension in OFFICE_PPTX_EXTENSIONS:
        pres = Presentation(path)
        text = []

        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)

        return "\n".join(text)

    elif extension in OFFICE_XLSX_EXTENSIONS:
        book = openpyxl.load_workbook(path, data_only=True)
        text = []

        for sheet in book.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        text.append(str(cell.value))

        return "\n".join(text)

    elif extension in MAC_EXTENSIONS:
        with zipfile.ZipFile(path, 'r') as z:
            text = []

            for name in z.namelist():
                if name.endswith((".xml", ".plist", ".rtf", ".txt")):
                    with z.open(name) as file:
                        content = file.read().decode("utf-8", errors="ignore")
                        text.append(content)

            return "\n".join(text)

    elif extension in EMAIL_EXTENSIONS:
        with open(path, encoding="utf-8", errors="ignore") as f:
            msg = email.message_from_file(f)
            parts = []

            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    payload = part.get_payload(decode=True)

                    if payload:
                        parts.append(payload.decode("utf-8", errors="ignore"))

            return "\n".join(parts)

    raise UnsupportedExtensionError(f"The provided file extension is not supported for metadata: .{extension}")


def translate(text: str) -> str:
    try:
        detected_lang = detect(text)

        if detected_lang == "en":
            return text

        else:
            translator = Translator()
            translation = translator.translate(text, dest="en")

            return translation.text

    except Exception as e:
        raise ProcessingError(f"Failed to translate the provided text: {str(e)}")


def sanitize(text: str) -> str:
    from nltk.tokenize import word_tokenize  # Lazy loading for faster startup
    from nltk.tokenize.treebank import TreebankWordDetokenizer

    # Normalize Unicode (NFKC handles compatibility + decomposition)
    text: str = unicodedata.normalize("NFKC", text)

    # Replace non-breaking spaces and other invisible chars with a normal space
    text = re.sub(r"[\u00A0\u200B-\u200D\uFEFF]", " ", text)

    # Remove control characters (except common whitespace chars)
    text = re.sub(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]", "", text)

    # Normalize whitespace
    text = re.sub(r"\s+", " ", text).strip()

    # Tokenize and detokenize (to normalize spacing around punctuation)
    tokens = word_tokenize(text)
    text = TreebankWordDetokenizer().detokenize(tokens)
    text = re.sub(r"\s+([?.!,;:])", r"\1", text)

    return text.lower()


def get_chunks(text: str, max_chunk_size: int = 500) -> list[str]:
    from nltk import download  # Lazy loading for faster startup
    from nltk.tokenize import sent_tokenize

    download("punkt_tab", quiet=True)

    chunks = []

    if _estimate_token_count(text) <= max_chunk_size:
        return [text.strip()]

    parts = sent_tokenize(text)

    for part in parts:
        part = part.strip()

        if not part:
            continue

        if _estimate_token_count(part) > max_chunk_size:
            chunks.extend(get_chunks(part, max_chunk_size))

        else:
            chunks.append(part)

    return chunks


def get_keywords(text: str) -> list[str]:
    cleaned = ''.join(' ' if unicodedata.category(c).startswith("P") else c for c in text)
    words = cleaned.split()

    seen = set()
    unique_words = []

    for word in words:
        if word not in seen:
            seen.add(word)
            unique_words.append(word)

    return unique_words


def get_text_summary(text: str, model: str = TEXT_SUMMARY_MODEL) -> str:
    from transformers import pipeline  # Lazy loading for faster startup

    summarizer = pipeline("summarization", model=model)
    size_ratio: float = 0.25
    summary = summarizer(text, max_length=int(size_ratio * _estimate_token_count(text)), do_sample=False)

    return summary[0]["summary_text"]


def get_text_embedding(text: str, model: str = TEXT_EMBED_MODEL) -> np.ndarray:
    # from sentence_transformers import (
    #     SentenceTransformer,  # Lazy loading for faster startup
    # )
    #
    # encoder = SentenceTransformer(
    #     model,
    #     tokenizer_kwargs={"padding_side": "left"},
    #     truncate_dim=48
    # )
    # embedding: np.ndarray = encoder.encode(text, convert_to_numpy=True, normalize_embeddings=True)

    from transformers import (  # Lazy loading for faster startup
        AutoModel,
        AutoTokenizer,
    )

    tokenizer = AutoTokenizer.from_pretrained(model)
    encoder = AutoModel.from_pretrained(model)
    inputs = tokenizer(text, padding=True, truncation=True, return_tensors="pt")

    outputs = encoder(**inputs)

    def mean_pooling(token_embeddings, mask):
        token_embeddings = token_embeddings.masked_fill(~mask[..., None].bool(), 0.)
        sentence_embeddings = token_embeddings.sum(dim=1) / mask.sum(dim=1)[..., None]

        return sentence_embeddings

    embeddings = mean_pooling(outputs[0], inputs["attention_mask"])
    embeddings = torch.nn.functional.normalize(embeddings[:, :16], p=2, dim=1)

    return quantize_embedding(embeddings.detach().numpy())


def text_to_bits(text: str) -> np.ndarray:
    encoded_string = text.encode("utf-8")
    hashed_data = hash(encoded_string)

    if hashed_data >= 0:
        binary = bin(hashed_data)[2:].zfill(32)

    else:
        binary = bin((1 << 32) + hashed_data)[2:].zfill(32)

    bit_list: list[int] = [int(bit) for bit in binary][:16]

    return np.array(bit_list)
